from __future__ import annotations

from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.util import Pt

from src.core.document import DocumentABC
from src.rendering.styles import TextStyle


class PptEngine(DocumentABC):

    def __init__(self, path: str | Path | None = None) -> None:
        super().__init__(path)
        self._prs: Presentation | None = None
        self._base_style: TextStyle = TextStyle.default_ppt()

    @property
    def prs(self) -> Presentation:
        if self._prs is None:
            raise RuntimeError('No presentation loaded. Call create() or open() first.')
        return self._prs

    def create(self) -> None:
        self._prs = Presentation()
        self._path = None
        self._base_style = TextStyle.default_ppt()

    def open(self, path: str | Path) -> None:
        self._path = Path(path)
        if not self._path.exists():
            raise FileNotFoundError(f'File not found: {self._path}')
        self._prs = Presentation(str(self._path))
        self._base_style = TextStyle.default_ppt()

    def save(self, path: str | Path | None = None) -> None:
        if path is not None:
            self._path = Path(path)
        if self._path is None:
            raise ValueError('No output path specified.')
        self._path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(self._path))

    def get_base_style(self) -> TextStyle:
        return self._base_style

    def add_text(
        self,
        text: str,
        bold: bool = False,
        italic: bool = False,
        font_name: str | None = None,
        font_size: int | None = None,
        color: str | None = None,
        alignment: str | None = None,
        text_style: TextStyle | None = None,
        **kwargs: Any,
    ) -> Any:
        import re

        inline = TextStyle(
            bold=bold or None,
            italic=italic or None,
            font_name=font_name,
            font_size=font_size,
            color=color,
            alignment=alignment,
        )
        final = self._base_style
        if text_style is not None:
            final = final.merge(text_style)
        final = final.merge(inline)

        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title

        if not title_shape:
            return slide

        tf = title_shape.text_frame
        has_math = '$$' in text or ('$' in text and re.search(r'\$[^$]+\$', text))

        if has_math:
            parts = re.split(r'(\$\$|\$)'.format(), text)
            p = tf.paragraphs[0]
            in_math = False
            math_display = False
            buffer = ''
            for part in parts:
                if part == '$$':
                    if in_math and math_display and buffer:
                        self._add_omath_to_paragraph(p, buffer)
                        buffer = ''
                    math_display = not math_display
                    in_math = math_display
                elif part == '$':
                    if in_math and not math_display and buffer:
                        self._add_omath_to_paragraph(p, buffer)
                        buffer = ''
                    math_display = False
                    in_math = not in_math
                elif in_math:
                    buffer += part
                else:
                    run = p.add_run()
                    run.text = part
                    self._apply_run_style(run, final)
            if buffer and in_math:
                self._add_omath_to_paragraph(p, buffer)
        else:
            tf.text = text
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    self._apply_run_style(run, final)
                if final.alignment:
                    from pptx.enum.text import PP_ALIGN
                    align_map = {
                        'left': PP_ALIGN.LEFT,
                        'center': PP_ALIGN.CENTER,
                        'right': PP_ALIGN.RIGHT,
                        'justify': PP_ALIGN.JUSTIFY,
                    }
                    paragraph.alignment = align_map.get(final.alignment, PP_ALIGN.LEFT)
        return slide

    def _apply_run_style(self, run: Any, style: TextStyle) -> None:
        if style.bold is not None and style.bold:
            run.font.bold = True
        if style.italic is not None and style.italic:
            run.font.italic = True
        if style.font_name is not None:
            run.font.name = style.font_name
        if style.font_size is not None:
            run.font.size = Pt(style.font_size)
        if style.color is not None:
            try:
                from pptx.dml.color import RGBColor
                run.font.color.rgb = RGBColor.from_string(style.color)
            except Exception:
                pass

    def add_styled_content(
        self,
        tokens: list[dict[str, Any]],
    ) -> Any:
        slide = self.add_slide()
        current_style = self._base_style

        for token in tokens:
            token_type = token.get('type', 'text')
            content = token.get('content', '')

            if token_type == 'text':
                if content.strip():
                    self._append_text_to_slide(slide, content, current_style)
            elif token_type == 'command':
                cmd = token.get('command', '')
                if cmd in ('newpage',):
                    slide = self.add_slide()
                    continue
                if cmd == 'heading':
                    h_content = token.get('content', '')
                    self._append_text_to_slide(slide, h_content, current_style, heading=True)
                    continue
                if cmd == 'math':
                    latex = token.get('latex', '')
                    self._add_omath_to_slide(slide, latex, current_style)
                    continue
                token_style = TextStyle.from_latex_token(token)
                merged = current_style.merge(token_style)
                inner = token.get('content', '')
                if inner:
                    self._append_text_to_slide(slide, inner, merged)

        return slide

    def _add_omath_to_slide(
        self, slide: Any, latex: str, style: TextStyle
    ) -> None:
        from pptx.util import Inches

        from src.rendering.math_omml import latex_to_omml

        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(1)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        omml = latex_to_omml(latex)
        if omml is not None:
            p._p.append(omml)

    def _add_omath_to_paragraph(self, paragraph: Any, latex: str) -> None:
        from src.rendering.math_omml import latex_to_omml
        omml = latex_to_omml(latex)
        if omml is not None:
            paragraph._p.append(omml)

    def _append_text_to_slide(
        self, slide: Any, text: str, style: TextStyle, heading: bool = False
    ) -> None:
        from pptx.util import Inches

        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(1)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if heading:
            p.level = 0
        run = p.add_run()
        run.text = text
        self._apply_run_style(run, style)

    def add_latex_content(self, text: str) -> Any:
        from src.rendering.latex_parser import parse_structured
        tokens = parse_structured(text)
        return self.add_styled_content(tokens)

    def replace_text(self, old: str, new: str, regex: bool = False) -> int:
        import re as _re
        count = 0
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if regex:
                                new_text = _re.sub(old, new, run.text)
                                if new_text != run.text:
                                    run.text = new_text
                                    count += 1
                            else:
                                if old in run.text:
                                    run.text = run.text.replace(old, new)
                                    count += 1
        return count

    def set_style(self, style_str: str) -> None:
        new_style = TextStyle.from_string(style_str)
        self._base_style = self._base_style.merge(new_style)

    def apply_style_to_all(self) -> None:
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            self._apply_run_style(run, self._base_style)

    def to_pdf(self, output_path: str | Path) -> None:
        self.save()
        from src.transform.pdf import ppt_to_pdf
        ppt_to_pdf(str(self._path), str(output_path))

    def add_slide(self, layout_index: int = 1) -> Any:
        slide_layout = self.prs.slide_layouts[layout_index]
        return self.prs.slides.add_slide(slide_layout)

    def apply_layout(self, slide_index: int, layout_spec: str) -> None:
        try:
            layout_idx = int(layout_spec)
            slide_layout = self.prs.slide_layouts[layout_idx]
        except ValueError:
            slide_layout = None
            for lo in self.prs.slide_layouts:
                if lo.name.lower() == layout_spec.lower():
                    slide_layout = lo
                    break
            if slide_layout is None:
                raise ValueError(
                    f'Layout "{layout_spec}" not found. '
                    f'Available: {[lo.name for lo in self.prs.slide_layouts]}'
                )
        if slide_index < len(self.prs.slides):
            slide = self.prs.slides[slide_index]
            ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            layout_part = slide_layout.part
            r_id = slide.part.relate_to(layout_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout')
            csld = slide.element.find(f'{{{ns}}}cSld')
            if csld is not None:
                layout_node = csld.find(f'{{{ns}}}sldLayout')
                if layout_node is not None:
                    csld.remove(layout_node)
                layout_node = etree.SubElement(csld, f'{{{ns}}}sldLayout')
                layout_node.set(f'{{{r_ns}}}id', r_id)

    def delete_slide(self, index: int) -> None:
        r_id = self.prs.slides._sldIdLst[index].attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
        self.prs.part.drop_rel(r_id)
        del self.prs.slides._sldIdLst[index]

    def move_slide(self, from_index: int, to_index: int) -> None:
        slides = list(self.prs.slides._sldIdLst)
        slide_entry = slides.pop(from_index)
        slides.insert(to_index, slide_entry)
        self.prs.slides._sldIdLst[:] = slides

    def get_metadata(self) -> dict[str, str | None]:
        props = self.prs.core_properties
        return {
            'author': props.author,
            'title': props.title,
            'subject': props.subject,
            'category': props.category,
            'keywords': props.keywords,
            'comments': props.comments,
        }

    def set_metadata(self, **kwargs: str) -> None:
        props = self.prs.core_properties
        for key, value in kwargs.items():
            key_lower = key.lower()
            if key_lower == 'author':
                props.author = value
            elif key_lower == 'title':
                props.title = value
            elif key_lower == 'subject':
                props.subject = value
            elif key_lower == 'category':
                props.category = value
            elif key_lower == 'keywords':
                props.keywords = value
            elif key_lower == 'comments':
                props.comments = value

    def add_comment(self, text: str, slide_index: int = 0) -> None:
        if slide_index >= len(self.prs.slides):
            raise ValueError(f'Slide index {slide_index} out of range.')
        slide = self.prs.slides[slide_index]
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        if tf.text:
            tf.text += '\n---\n'
        tf.text += text

    def add_notes(self, slide_index: int, text: str) -> None:
        if slide_index >= len(self.prs.slides):
            raise ValueError(f'Slide index {slide_index} out of range.')
        notes_slide = self.prs.slides[slide_index].notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = text

    def to_images(self, output_dir: str | Path) -> list[Path]:
        import shutil
        import subprocess
        from pathlib import Path

        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        lo_bin = None
        lo_paths = [
            'libreoffice', 'soffice',
            '/usr/bin/libreoffice', '/usr/bin/soffice',
            r'C:\Program Files\LibreOffice\program\soffice.exe',
        ]
        for p in lo_paths:
            if shutil.which(p):
                lo_bin = p
                break

        if not lo_bin:
            raise RuntimeError(
                'LibreOffice is required for slide image export. '
                'Install from https://www.libreoffice.org/download/'
            )

        self.save()
        subprocess.run(
            [lo_bin, '--headless', '--convert-to', 'png',
             '--outdir', str(output_dir), str(self._path)],
            check=True, capture_output=True,
        )

        return sorted(output_dir.glob('*.png'))

    def set_transition(self, transition_type: str, slide_index: int | None = None) -> None:
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        valid_transitions = {
            'fade', 'push', 'wipe', 'cover', 'uncover', 'dissolve',
            'random', 'split', 'strips', 'blinds', 'checker', 'comb',
            'zoom', 'glitter', 'vortex', 'ripple', 'honeycomb',
        }
        ttype = transition_type.lower()
        if ttype not in valid_transitions:
            raise ValueError(
                f'Unsupported transition: {transition_type}. '
                f'Use one of: {", ".join(sorted(valid_transitions))}'
            )
        slides = [self.prs.slides[slide_index]] if slide_index is not None \
            else list(self.prs.slides)
        for slide in slides:
            existing = slide.element.findall(f'{{{ns}}}transition')
            for el in existing:
                slide.element.remove(el)
            trans_el = etree.SubElement(slide.element, f'{{{ns}}}transition')
            etree.SubElement(trans_el, f'{{{ns}}}{ttype}')

    def set_protection(self, password: str) -> None:
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        pres_elem = self.prs.part._element
        existing = pres_elem.findall(f'{{{ns}}}modifyVerifier')
        for el in existing:
            pres_elem.remove(el)
        verifier = etree.SubElement(pres_elem, f'{{{ns}}}modifyVerifier')
        verifier.set('cryptProviderType', 'rsaAES')
        verifier.set('cryptAlgorithmClass', 'hash')
        verifier.set('cryptAlgorithmType', 'typeAny')
        verifier.set('cryptAlgorithmSid', '14')
        verifier.set('cryptSpinCount', '100000')
        verifier.set('hashData', password)
        verifier.set('saltData', password)

    def unprotect(self) -> None:
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        pres_elem = self.prs.part._element
        existing = pres_elem.findall(f'{{{ns}}}modifyVerifier')
        for el in existing:
            pres_elem.remove(el)
