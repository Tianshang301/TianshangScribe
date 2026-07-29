"""Sub-agent: Create 3 Office documents via TianshangScribe MCP server, then verify.

Protocol: stdio JSON-RPC. Each document = separate subprocess invocation.
"""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from threading import Thread

MCP_DIR = Path(r"F:\Projects\Project20")
OUTPUT_DIR = MCP_DIR / "mcp" / "subagent_output"


# ─── MCP client helpers ────────────────────────────────────────────────────

def _start_mcp_server(stderr_capture: bool = False) -> subprocess.Popen:
    proc = subprocess.Popen(
        [sys.executable, "-m", "mcp.server"],
        stdin=subprocess.PIPE,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE if stderr_capture else subprocess.DEVNULL,
        cwd=str(MCP_DIR),
        text=True,
        bufsize=1,
    )
    resp = _send_and_read(proc, {
        "jsonrpc": "2.0", "id": 1, "method": "initialize",
        "params": {"protocolVersion": "2024-11-05", "capabilities": {}}
    })
    if "error" in resp:
        raise RuntimeError(f"Initialize failed: {resp}")
    proc.stdin.write(
        json.dumps({"jsonrpc": "2.0", "method": "notifications/initialized"}) + "\n"
    )
    proc.stdin.flush()
    return proc


def _send_and_read(proc: subprocess.Popen, request: dict) -> dict:
    line = json.dumps(request, ensure_ascii=False)
    proc.stdin.write(line + "\n")
    proc.stdin.flush()
    resp_line = proc.stdout.readline()
    if not resp_line:
        raise RuntimeError("MCP server closed stdout unexpectedly")
    return json.loads(resp_line)


def _call_tool(proc: subprocess.Popen, name: str, args: dict) -> dict:
    resp = _send_and_read(proc, {
        "jsonrpc": "2.0", "id": 99, "method": "tools/call",
        "params": {"name": name, "arguments": args}
    })
    if "error" in resp:
        return {"error": resp["error"]}
    try:
        inner = json.loads(resp["result"]["content"][0]["text"])
    except (KeyError, IndexError, json.JSONDecodeError) as e:
        return {"error": str(e), "raw": resp}
    # Unwrap success envelope: {"success": true, "data": {...}}
    if inner.get("success") and "data" in inner:
        return inner["data"]
    return inner


def _run_doc(args: dict, label: str) -> dict:
    """Start a fresh MCP server, make one call, return result dict."""
    stderr_lines: list[str] = []

    proc = _start_mcp_server(stderr_capture=True)

    def _drain():
        try:
            while True:
                line = proc.stderr.readline()
                if not line:
                    break
                stderr_lines.append(line.rstrip("\n"))
        except Exception:
            pass

    t = Thread(target=_drain, daemon=True)
    t.start()

    try:
        result = _call_tool(proc, args["tool"], args["params"])
    finally:
        try:
            proc.stdin.close()
            proc.wait(timeout=5)
        except Exception:
            proc.kill()
        t.join(timeout=1)

    path = result.get("output_path", "") if isinstance(result, dict) else str(result)
    err_info = ""
    if isinstance(result, dict) and "error" in result:
        err_info = f" | ERROR: {result['error']}"
        if stderr_lines:
            err_info += f" | stderr: {' // '.join(stderr_lines)}"
    print(f"  [{label}] {path}{err_info}")
    return result


# ─── Document definitions ──────────────────────────────────────────────────

def doc_word(tmpdir: str) -> dict:
    path = os.path.join(tmpdir, "annual_review.docx")
    return {
        "tool": "create_office_document",
        "params": {
            "format": "docx",
            "output_path": path,
            "style": "font=Calibri,size=12",
            "metadata": {"title": "2026 Annual Technology Review", "author": "TianshangScribe Agent", "subject": "Annual Review"},
            "content": [
                {"type": "heading", "level": 1, "text": "2026 Annual Technology Review"},
                {"type": "heading", "level": 2, "text": "Executive Summary"},
                {"type": "paragraph", "text": (
                    r"This report summarizes key technology trends and organizational "
                    r"achievements in the first half of 2026. \bfseries{Highlights:} "
                    r"AI adoption grew \color{0000FF}{340%}, cloud migration "
                    r"\itshape{completed 2 months ahead of schedule}, and our new "
                    r"MCP-based document automation pipeline reduced manual report "
                    r"generation time by \color{FF0000}{92%}."
                )},
                {"type": "heading", "level": 2, "text": "Financial Performance"},
                {"type": "formula", "text": r"\sum_{Q1}^{Q2} R_i = 12.5\text{M} + 14.8\text{M} = 27.3\text{M}"},
                {"type": "paragraph", "text": "Revenue trajectory is strong. The compound monthly growth rate (CMGR) follows:"},
                {"type": "formula", "text": r"\text{CMGR} = \left(\frac{V_f}{V_i}\right)^{\frac{1}{n}} - 1"},
                {"type": "heading", "level": 2, "text": "Team Expansion"},
                {"type": "table", "rows": [
                    ["Department", "Q1 Headcount", "Q2 Headcount", "Growth"],
                    ["Engineering", "45", "62", "+37.8%"],
                    ["AI/ML", "12", "28", "+133%"],
                    ["Product", "18", "22", "+22.2%"],
                    ["DevOps", "8", "15", "+87.5%"],
                ]},
                {"type": "heading", "level": 2, "text": "Key Metrics Dashboard"},
                {"type": "paragraph", "text": (
                    r"\bfseries{Uptime:} 99.97% "
                    r"\centering{\bfseries{Deployment Frequency: 47/day}} "
                    r"\raggedright{\itshape{Customer Satisfaction: 4.8/5.0}}"
                )},
                {"type": "page_break"},
                {"type": "heading", "level": 2, "text": "Risk Assessment"},
                {"type": "paragraph", "text": (
                    r"\bfseries{Critical:} \color{FF0000}{Supply chain disruption risk "
                    r"remains high.} Mitigation strategies include vendor diversification "
                    r"and strategic stockpiling. \bfseries{Moderate:} Talent acquisition "
                    r"competition in AI/ML space may impact hiring timelines."
                )},
                {"type": "heading", "level": 2, "text": "Recommendations"},
                {"type": "paragraph", "text": (
                    r"1. Increase AI/ML budget by \color{0000FF}{25%} in Q3."
                    r"\newpage 2. Launch MCP integration for all document workflows "
                    r"by September.\newpage 3. Establish vendor redundancy for "
                    r"critical components."
                )},
            ],
        },
    }


def doc_excel(tmpdir: str) -> dict:
    path = os.path.join(tmpdir, "sales_dashboard.xlsx")
    return {
        "tool": "create_office_document",
        "params": {
            "format": "xlsx",
            "output_path": path,
            "style": "font=Calibri,size=11",
            "content": [
                {"type": "paragraph", "text": "Product"},
                {"type": "paragraph", "text": "Widget Pro"},
                {"type": "paragraph", "text": "Widget Lite"},
                {"type": "paragraph", "text": "SuperGadget"},
                {"type": "paragraph", "text": "MegaGizmo"},
                {"type": "paragraph", "text": "NanoThing"},
                {"type": "paragraph", "text": "UltraDoodad"},
                {"type": "paragraph", "text": ""},
                {"type": "paragraph", "text": "Total Revenue"},
            ],
        },
    }


def doc_ppt(tmpdir: str) -> dict:
    path = os.path.join(tmpdir, "mcp_strategy.pptx")
    return {
        "tool": "create_office_document",
        "params": {
            "format": "pptx",
            "output_path": path,
            "content": [
                {"type": "heading", "level": 1, "text": "MCP Integration Strategy"},
                {"type": "paragraph", "text": r"\bfseries{TianshangScribe + MCP = Agent-Native Document Automation}"},
                {"type": "heading", "level": 2, "text": "Problem Statement"},
                {"type": "paragraph", "text": (
                    "Current document workflows require manual intervention. Agents "
                    "cannot reliably create, edit, or convert Office documents. This "
                    "creates a bottleneck in automated reporting pipelines."
                )},
                {"type": "heading", "level": 2, "text": "Solution Architecture"},
                {"type": "paragraph", "text": (
                    r"TianshangScribe MCP Server bridges the gap between AI Agents "
                    r"and Office formats. Agents use JSON-RPC via stdio to create, edit, "
                    r"and convert documents \u2014 with \bfseries{zero external dependencies.}"
                )},
                {"type": "formula", "text": r"\text{Agent} \xrightarrow{\text{JSON-RPC}} \text{MCP Server} \xrightarrow{\text{OOML}} \text{Office Document}"},
                {"type": "heading", "level": 2, "text": "Integration Points"},
                {"type": "paragraph", "text": (
                    r"\bfseries{Claude Code:} Direct MCP client integration via config file."
                    r"\newpage \bfseries{Cursor IDE:} MCP tools exposed to AI assistant."
                    r"\newpage \bfseries{Custom Agents:} stdio JSON-RPC, easy to embed."
                    r"\newpage \bfseries{Coze/Dify:} Planned SSE support for cloud agents."
                )},
                {"type": "heading", "level": 2, "text": "Performance Metrics"},
                {"type": "paragraph", "text": (
                    r"\centering{Document generation time: \color{0000FF}{<2 seconds}} "
                    r"\centering{Supported formats: \bfseries{8}} "
                    r"\centering{Agent success rate: \color{0000FF}{99.7%}}"
                )},
                {"type": "heading", "level": 2, "text": "Next Steps"},
                {"type": "paragraph", "text": (
                    r"Q3: Launch SSE transport, Coze plugin marketplace."
                    r"\newpage Q4: Multi-language template support, real-time collaboration API."
                )},
            ],
        },
    }


# ─── Verification ──────────────────────────────────────────────────────────

def verify_word(path: str) -> dict:
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(path)
    paragraphs = doc.paragraphs
    tables = doc.tables

    bold_count = 0
    italic_count = 0
    omml_count = 0
    for p in paragraphs:
        for run in p.runs:
            if run.bold:
                bold_count += 1
            if run.italic:
                italic_count += 1
        omml_count += len(p._p.findall(qn("m:oMath")))
        omml_count += len(p._p.findall(qn("m:oMathPara")))

    return {
        "paragraphs": len(paragraphs),
        "tables": len(tables),
        "runs_bold": bold_count,
        "runs_italic": italic_count,
        "omml_elements": omml_count,
        "sections": len(doc.sections),
        "title": doc.core_properties.title,
        "author": doc.core_properties.author,
    }


def verify_excel(path: str) -> dict:
    from openpyxl import load_workbook

    wb = load_workbook(path)
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    return {
        "sheets": len(wb.sheetnames),
        "sheet_names": wb.sheetnames,
        "max_row": ws.max_row,
        "max_col": ws.max_column,
        "sample_rows": [(i + 1, list(r)) for i, r in enumerate(rows[:12])],
    }


def verify_ppt(path: str) -> dict:
    from pptx import Presentation

    prs = Presentation(path)
    slides_info = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t[:80])
        slides_info.append({"index": i + 1, "shapes": len(slide.shapes), "texts": texts[:3]})
    return {"slide_count": len(prs.slides), "slides": slides_info}


# ─── Main ──────────────────────────────────────────────────────────────────

def main():
    os.makedirs(str(OUTPUT_DIR), exist_ok=True)
    tmpdir = tempfile.mkdtemp(prefix="scribe_subagent_")
    print(f"Working dir: {tmpdir}")
    print(f"Output dir : {OUTPUT_DIR}")
    print()

    results: dict[str, dict] = {}

    # ── Document 1: Word ──
    print("=" * 60)
    print("  DOCUMENT 1: Word")
    print("=" * 60)
    r = _run_doc(doc_word(tmpdir), "Word")
    results["word"] = r

    # ── Document 2: Excel ──
    print()
    print("=" * 60)
    print("  DOCUMENT 2: Excel")
    print("=" * 60)
    r = _run_doc(doc_excel(tmpdir), "Excel")
    results["excel"] = r

    # ── Fix Excel columns 2 & 3 ──
    excel_path = results.get("excel", {}).get("output_path", "")
    if excel_path and os.path.exists(excel_path):
        print(f"  [Excel edit] Adding columns B and C via openpyxl...")
        try:
            from openpyxl import load_workbook
            from openpyxl.styles import Font

            wb = load_workbook(excel_path)
            ws = wb.active
            col2_data = ["Price", "45.99", "19.99", "89.50", "34.00", "12.75", "67.30"]
            col3_data = ["Units", "1,200", "3,450", "890", "2,100", "5,600", "780"]
            font_calibri = Font(name="Calibri", size=11)
            for idx, (c2, c3) in enumerate(zip(col2_data, col3_data), start=1):
                c = ws.cell(row=idx, column=2, value=c2)
                c.font = font_calibri
                c = ws.cell(row=idx, column=3, value=c3)
                c.font = font_calibri
            wb.save(excel_path)
            print(f"  [Excel edit] Done — 7 rows x 3 columns written.")
        except Exception as exc:
            print(f"  [Excel edit] ERROR: {exc}")

    # ── Document 3: PPT ──
    print()
    print("=" * 60)
    print("  DOCUMENT 3: PPT")
    print("=" * 60)
    r = _run_doc(doc_ppt(tmpdir), "PPT")
    results["ppt"] = r

    # ── Verification ──
    print()
    print("=" * 60)
    print("  VERIFICATION")
    print("=" * 60)

    word_path = results.get("word", {}).get("output_path", "")
    excel_path = results.get("excel", {}).get("output_path", "")
    ppt_path = results.get("ppt", {}).get("output_path", "")

    for label, path, verifier in [
        ("Word", word_path, verify_word),
        ("Excel", excel_path, verify_excel),
        ("PPT", ppt_path, verify_ppt),
    ]:
        print(f"\n-- {label}: {path} --")
        if not path:
            print(f"   FAILED: {results.get(label.lower(), {})}")
            continue
        if not os.path.exists(path):
            print(f"   FILE NOT FOUND on disk")
            continue
        try:
            v = verifier(path)
            v["file_size_bytes"] = os.path.getsize(path)
            for k, val in v.items():
                print(f"   {k}: {val}")
        except Exception as exc:
            import traceback
            print(f"   VERIFICATION ERROR: {exc}")
            traceback.print_exc()

    # ── Copy to output dir ──
    print()
    print("=" * 60)
    print("  COPY TO OUTPUT DIRECTORY")
    print("=" * 60)

    for label, path in [("Word", word_path), ("Excel", excel_path), ("PPT", ppt_path)]:
        if path and os.path.exists(path):
            dst = str(OUTPUT_DIR / os.path.basename(path))
            shutil.copy2(path, dst)
            print(f"  {label}: {dst}  ({os.path.getsize(dst):,} bytes)")
        else:
            print(f"  {label}: SKIPPED — source missing: {path}")

    print()
    print("=" * 60)
    print("  COMPLETE")
    print("=" * 60)


if __name__ == "__main__":
    main()
