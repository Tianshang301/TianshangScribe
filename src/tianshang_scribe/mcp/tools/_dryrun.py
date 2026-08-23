"""Pre-flight validation shared by the edit/create tools' ``dry_run`` mode.

Given a planned operation list, :func:`build_edit_plan` opens the target file
read-only for context (sheet names / slide count), then checks each operation's
references structurally: cell references, ranges, sheet existence, and PPT
slide-index bounds. The plan reports per-operation findings with refined error
codes so Agents can fix inputs BEFORE paying for a real write.

Deliberately out of scope: Excel formula syntax validation. A shallow grammar
check gives false confidence and a real one is half a parser — formulas are
passed to the engine verbatim.

Live-error refinement lives here too: :func:`classify_engine_error` maps engine
``ValueError``/``IndexError`` messages onto the refined codes so in-flight
failures carry the same taxonomy as dry-run findings.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from openpyxl.utils.exceptions import CellCoordinatesException

from tianshang_scribe.mcp.errors import McpErrorCode

_MAX_COL = 16384  # XFD
_MAX_ROW = 1_048_576

_CELL_RE = re.compile(r'^([A-Z]{1,3})([0-9]+)$')
_RANGE_RE = re.compile(r'^([A-Z]{1,3})([0-9]+):([A-Z]{1,3})([0-9]+)$')
_ROW_RANGE_RE = re.compile(r'^([0-9]+):([0-9]+)$')
_COL_RANGE_RE = re.compile(r'^([A-Z]{1,3}):([A-Z]{1,3})$')

#: Engine message patterns -> (refined code, offending field). First match wins.
_ERROR_PATTERNS: tuple[tuple[re.Pattern[str], int, str], ...] = (
    (
        re.compile(r'slide_index out of range|slide index \d+ out of range', re.IGNORECASE),
        McpErrorCode.PPT_INVALID_SLIDE_INDEX,
        'slide_index',
    ),
    (re.compile(r'Sheet not found'), McpErrorCode.EXCEL_SHEET_NOT_FOUND, 'sheet_name'),
    (
        re.compile(r'Invalid (?:row|column|cell) range|Invalid print area'),
        McpErrorCode.EXCEL_INVALID_RANGE,
        'range',
    ),
    (re.compile(r'is_formula'), McpErrorCode.INVALID_PARAMETER, 'value'),
)


def _col_to_num(letters: str) -> int:
    num = 0
    for ch in letters:
        num = num * 26 + (ord(ch) - ord('A') + 1)
    return num


def _cell_problem(ref: Any) -> str | None:
    """Return a problem description for ``ref``, or None when it is a valid cell."""
    if not isinstance(ref, str):
        return f'cell reference must be a string, got {type(ref).__name__}'
    match = _CELL_RE.match(ref.strip().upper())
    if match is None:
        return f'{ref!r} is not an A1-style reference like "B2"'
    col, row = _col_to_num(match.group(1)), int(match.group(2))
    if not 1 <= col <= _MAX_COL:
        return f'column {match.group(1)!r} exceeds XFD'
    if not 1 <= row <= _MAX_ROW:
        return f'row {row} outside 1-{_MAX_ROW}'
    return None


def _range_problem(rng: Any) -> str | None:
    """Return a problem description for ``rng``, or None when structurally valid."""
    if not isinstance(rng, str):
        return f'range must be a string, got {type(rng).__name__}'
    text = rng.strip().upper()
    if ':' not in text:
        return _cell_problem(text)
    match = _RANGE_RE.match(text)
    if match is None:
        return f'{rng!r} is not a "START:END" range like "A1:C10"'
    start_col, start_row = _col_to_num(match.group(1)), int(match.group(2))
    end_col, end_row = _col_to_num(match.group(3)), int(match.group(4))
    if max(start_col, end_col) > _MAX_COL or min(start_row, end_row) < 1:
        return f'range {rng!r} exceeds sheet bounds'
    if max(start_row, end_row) > _MAX_ROW:
        return f'range {rng!r} exceeds sheet bounds'
    return None


def validate_cell_ref(ref: str) -> tuple[bool, str]:
    """Validate an A1-style cell reference; return ``(ok, detail)``."""
    problem = _cell_problem(ref)
    return (problem is None, problem or '')


def validate_range(rng: str) -> tuple[bool, str]:
    """Validate an A1:B2-style range (single cells allowed); return ``(ok, detail)``."""
    problem = _range_problem(rng)
    return (problem is None, problem or '')


def estimate_range_cells(rng: str) -> int:
    """Cell count covered by an A1:B2 range; 0 when unparsable."""
    match = _RANGE_RE.match(rng.strip().upper())
    if match is None:
        return 1 if _CELL_RE.match(rng.strip().upper()) else 0
    cols = abs(_col_to_num(match.group(3)) - _col_to_num(match.group(1))) + 1
    rows = abs(int(match.group(4)) - int(match.group(2))) + 1
    return cols * rows


def classify_engine_error(exc: Exception) -> tuple[int, str] | None:
    """Map an engine error onto ``(refined_code, field)``, or None when unmatched.

    Unmatched errors keep their legacy mapping at the call site (ValueError ->
    DOCUMENT_LOCKED, IndexError -> INTERNAL_ERROR).
    """
    if isinstance(exc, CellCoordinatesException):
        return (McpErrorCode.EXCEL_INVALID_CELL_REF, 'cell')
    message = str(exc)
    for pattern, code, field in _ERROR_PATTERNS:
        if pattern.search(message):
            return (code, field)
    return None


def _finding(index: int, action: Any, field: str, code: int, detail: str) -> dict[str, Any]:
    return {
        'index': index,
        'action': action,
        'ok': False,
        'field': field,
        'error_code': code,
        'detail': detail,
    }


# Operations whose ``range`` must be a row range ("2:5"), a column range
# ("B:D"), or either form (ungroup).
_ROW_RANGE_ACTIONS = frozenset({'group_rows'})
_COL_RANGE_ACTIONS = frozenset({'group_columns'})

_CELL_FIELD_ACTIONS = {
    'write_cell': 'cell',
    'set_formula': 'cell',
}

_RANGE_SPEC_ACTIONS = {  # "RANGE=..." spec fields
    'number_format': 'number_format',
    'conditional_format': 'conditional_format',
    'data_validation': 'data_validation',
}

_A1_RANGE_ACTIONS = frozenset({'set_print_area', 'sort', 'set_range_style'})

#: Every action any edit surface can dispatch: the legacy shared dispatcher plus
#: the Excel/PPT wrapper-local edge actions (PLAN.md D-7). Dry-run flags
#: anything outside this set exactly as live dispatch would (1006).
KNOWN_ACTIONS = frozenset(
    {
        'replace',
        'delete',
        'modify',
        'style',
        'add',
        'clear',
        'write_cell',
        'set_formula',
        'freeze_panes',
        'add_chart',
        'conditional_format',
        'data_validation',
        'add_table',
        'add_picture',
        'add_shape',
        'sort',
        'add_sheet',
        'set_range_style',
        'number_format',
        'add_slide',
        'apply_layout',
        'set_transition',
        'add_notes',
        'group_rows',
        'group_columns',
        'ungroup',
        'add_media',
        'set_tab_color',
        'set_print_area',
        'set_page_setup',
        'apply_theme',
        'set_master_options',
    }
)


def _validate_op(op: dict[str, Any], index: int, ctx: _EditContext) -> dict[str, Any]:
    """Check one operation against ``ctx``; return a finding dict."""
    action = op.get('action')

    def ok() -> dict[str, Any]:
        return {'index': index, 'action': action, 'ok': True}

    if action not in KNOWN_ACTIONS:
        return _finding(
            index,
            action,
            'action',
            McpErrorCode.INVALID_PARAMETER,
            f'Unknown edit action: {action!r}.',
        )

    if action == 'add_sheet':
        name = op.get('sheet_name')
        if not name:
            return _finding(
                index,
                action,
                'sheet_name',
                McpErrorCode.INVALID_PARAMETER,
                'add_sheet requires sheet_name.',
            )
        ctx.add_sheet(name)
        return ok()

    if action == 'add_slide':
        # Extends the deck; later ops may target the new running index.
        if ctx.slide_count is not None:
            ctx.slide_count += 1
        return ok()

    # Sheet existence (Excel ops carrying sheet_name against known sheets).
    sheet = op.get('sheet_name')
    if sheet is not None and ctx.sheet_names is not None and sheet not in ctx.sheet_names:
        return _finding(
            index,
            action,
            'sheet_name',
            McpErrorCode.EXCEL_SHEET_NOT_FOUND,
            f'Sheet {sheet!r} does not exist; known: {", ".join(ctx.sheet_names[:8])}.',
        )

    # Slide bounds (running count accounts for add_slide earlier in the list).
    idx = op.get('slide_index')
    if (
        idx is not None
        and ctx.slide_count is not None
        and (not isinstance(idx, int) or isinstance(idx, bool) or not 0 <= idx < ctx.slide_count)
    ):
        return _finding(
            index,
            action,
            'slide_index',
            McpErrorCode.PPT_INVALID_SLIDE_INDEX,
            f'slide_index {idx!r} outside 0-{ctx.slide_count - 1}.',
        )

    # Single-cell targets.
    cell_field = _CELL_FIELD_ACTIONS.get(str(action))
    if cell_field:
        cell = op.get(cell_field)
        if cell is None:
            return _finding(
                index,
                action,
                cell_field,
                McpErrorCode.INVALID_PARAMETER,
                f'{action} requires {cell_field!r}.',
            )
        problem = _cell_problem(cell)
        if problem:
            return _finding(
                index,
                action,
                cell_field,
                McpErrorCode.EXCEL_INVALID_CELL_REF,
                problem,
            )
        return ok()

    if action == 'freeze_panes':
        anchor = op.get('range')
        if anchor is None:
            return _finding(
                index,
                action,
                'range',
                McpErrorCode.INVALID_PARAMETER,
                'freeze_panes requires range (anchor cell).',
            )
        problem = _cell_problem(anchor)
        if problem:
            return _finding(
                index,
                action,
                'range',
                McpErrorCode.EXCEL_INVALID_CELL_REF,
                problem,
            )
        return ok()

    if action in _ROW_RANGE_ACTIONS or action in _COL_RANGE_ACTIONS:
        rng = op.get('range')
        pattern = _ROW_RANGE_RE if action in _ROW_RANGE_ACTIONS else _COL_RANGE_RE
        if rng is None or not pattern.match(str(rng).strip().upper()):
            expected = '"2:5"' if action in _ROW_RANGE_ACTIONS else '"B:D"'
            return _finding(
                index,
                action,
                'range',
                McpErrorCode.EXCEL_INVALID_RANGE,
                f'grouping expects a range like {expected}, got {rng!r}.',
            )
        return ok()

    if action == 'ungroup':
        rng = op.get('range')
        axis = op.get('axis') or 'rows'
        pattern = _ROW_RANGE_RE if axis == 'rows' else _COL_RANGE_RE
        if rng is None or not pattern.match(str(rng).strip().upper()):
            expected = '"2:5"' if axis == 'rows' else '"B:D"'
            return _finding(
                index,
                action,
                'range',
                McpErrorCode.EXCEL_INVALID_RANGE,
                f'ungroup(axis={axis}) expects a range like {expected}, got {rng!r}.',
            )
        return ok()

    # A1:B2-style ranges and "RANGE=..." specs.
    spec_key = _RANGE_SPEC_ACTIONS.get(str(action))
    if spec_key is not None and op.get(spec_key) is not None:
        rng = str(op[spec_key]).partition('=')[0].strip()
        field = spec_key
    elif op.get('chart_data_range') is not None and action == 'add_chart':
        rng = str(op['chart_data_range']).strip()
        field = 'range'
    elif action in _A1_RANGE_ACTIONS:
        value = op.get('range')
        rng = value.strip() if isinstance(value, str) else ''
        field = 'range'
    else:
        return ok()
    problem = _range_problem(rng)
    if problem:
        return _finding(index, action, field, McpErrorCode.EXCEL_INVALID_RANGE, problem)
    return ok()


class _EditContext:
    """Mutable validation context: known sheets and running slide count."""

    def __init__(self, sheet_names: list[str] | None, slide_count: int | None) -> None:
        self.sheet_names = sheet_names
        self.slide_count = slide_count

    def add_sheet(self, name: str) -> None:
        if self.sheet_names is not None:
            self.sheet_names.append(name)


def _load_context(input_path: str) -> _EditContext:
    """Open the target read-only for sheet names / slide count; None on failure."""
    suffix = Path(input_path).suffix.lower()
    try:
        if suffix in ('.xlsx', '.xlsm'):
            from openpyxl import load_workbook

            wb = load_workbook(input_path, read_only=True, data_only=True)
            try:
                return _EditContext(list(wb.sheetnames), None)
            finally:
                wb.close()
        if suffix == '.pptx':
            from pptx import Presentation

            return _EditContext(None, len(Presentation(input_path).slides))
    except Exception:  # noqa: S110  # context is best-effort; structural checks still run
        pass
    return _EditContext(None, None)


def build_edit_plan(input_path: str, ops_list: list[dict[str, Any]]) -> dict[str, Any]:
    """Validate a planned operation list; return dry-run enrichment keys."""
    ctx = _load_context(input_path)
    validations: list[dict[str, Any]] = []
    impacted = 0
    for i, op in enumerate(ops_list):
        validations.append(_validate_op(op, i, ctx))
        impacted += _estimate_op_cells(op)
    return {
        'validations': validations,
        'all_valid': all(v['ok'] for v in validations),
        'estimated_impacted_cells': impacted,
    }


def _estimate_op_cells(op: dict[str, Any]) -> int:
    """Rough per-op impact estimate in cells; best-effort, defaults to 1."""
    action = op.get('action')
    if action in ('write_cell', 'set_formula', 'freeze_panes', 'add_sheet'):
        return 1
    if action == 'group_rows':
        match = _ROW_RANGE_RE.match(str(op.get('range') or '').strip())
        if match:
            return abs(int(match.group(2)) - int(match.group(1))) + 1
        return 0
    if action == 'group_columns':
        match = _COL_RANGE_RE.match(str(op.get('range') or '').strip().upper())
        if match:
            return abs(_col_to_num(match.group(2)) - _col_to_num(match.group(1))) + 1
        return 0
    if action == 'ungroup':
        return 0  # clears outline state, does not touch cell contents
    spec = op.get('number_format') or op.get('conditional_format') or op.get('data_validation')
    if isinstance(spec, str) and '=' in spec:
        return estimate_range_cells(spec.partition('=')[0].strip())
    rng = op.get('range')
    if isinstance(rng, str) and (':' in rng or _CELL_RE.match(rng.strip().upper())):
        return estimate_range_cells(rng)
    rows = op.get('rows')
    if isinstance(rows, list) and rows:
        width = len(rows[0]) if isinstance(rows[0], list) else 1
        return len(rows) * width
    chart_range = op.get('chart_data_range')
    if isinstance(chart_range, str):
        return estimate_range_cells(chart_range)
    return 1


def build_content_plan(blocks: list[dict[str, Any]]) -> dict[str, Any]:
    """Light-weight dry-run validation for ``create_office_document`` blocks.

    Only structural checks run here — the target file does not exist yet, so
    there are no sheet/slide contexts to resolve against.
    """
    validations: list[dict[str, Any]] = []
    impacted = 0
    for i, block in enumerate(blocks):
        block_type = block.get('type', 'paragraph')
        problems: list[dict[str, Any]] = []
        cell = block.get('cell')
        if cell is not None:
            problem = _cell_problem(cell)
            if problem:
                problems.append(
                    _finding(i, block_type, 'cell', McpErrorCode.EXCEL_INVALID_CELL_REF, problem)
                )
        for key in ('freeze', 'chart_data_range'):
            value = block.get(key)
            if value is not None:
                problem = _cell_problem(value) if key == 'freeze' else _range_problem(value)
                code = (
                    McpErrorCode.EXCEL_INVALID_CELL_REF
                    if key == 'freeze'
                    else McpErrorCode.EXCEL_INVALID_RANGE
                )
                if problem:
                    problems.append(_finding(i, block_type, key, code, problem))
        spec = block.get('number_format')
        if isinstance(spec, str) and '=' in spec:
            problem = _range_problem(spec.partition('=')[0].strip())
            if problem:
                problems.append(
                    _finding(
                        i,
                        block_type,
                        'number_format',
                        McpErrorCode.EXCEL_INVALID_RANGE,
                        problem,
                    )
                )
        if problems:
            validations.extend(problems)
        else:
            validations.append({'index': i, 'action': block_type, 'ok': True})
        rows = block.get('rows')
        if isinstance(rows, list) and rows:
            width = len(rows[0]) if isinstance(rows[0], list) and rows[0] else 1
            impacted += len(rows) * width
    return {
        'validations': validations,
        'all_valid': all(v['ok'] for v in validations),
        'estimated_impacted_cells': impacted,
    }
