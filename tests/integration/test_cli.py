"""CLI end-to-end tests.

The one-shot command is a plain Typer command, so options may appear before
or after the positional ``input_file``. ``open`` is dispatched to a separate
interactive subcommand app.
"""

from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from typer.testing import CliRunner

from src.cli.main import app, open_app
from src.core.word_engine import WordEngine

runner = CliRunner()


def _run(*args: str) -> tuple[int, str]:
    result = runner.invoke(app, list(args))
    return result.exit_code, result.stdout


class TestExtractCli:
    def test_extract_text(self, tmp_path: Path) -> None:
        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.add_text('Hello extract')
        e.save(str(docx))
        code, out = _run('--extract', 'text', str(docx))
        assert code == 0
        assert 'Hello extract' in out

    def test_extract_metadata(self, tmp_path: Path) -> None:
        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.add_text('x')
        e.save(str(docx))
        code, out = _run('--extract', 'metadata', str(docx))
        assert code == 0
        assert '"author"' in out

    def test_extract_tables(self, tmp_path: Path) -> None:
        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.add_table_data([['A', 'B'], ['1', '2']])
        e.save(str(docx))
        code, out = _run('--extract', 'tables', str(docx))
        assert code == 0
        assert 'A,B' in out
        assert '1,2' in out

    def test_extract_structure(self, tmp_path: Path) -> None:
        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.add_text('x')
        e.save(str(docx))
        code, out = _run('--extract', 'structure', str(docx))
        assert code == 0
        assert '"paragraphs"' in out


class TestAddTableCli:
    def test_add_table_inline(self, tmp_path: Path) -> None:
        out = tmp_path / 't.docx'
        code, _ = _run(
            '--force',
            '--create',
            '--word',
            '--output',
            str(out),
            '--add',
            'Hello',
            '--add-table',
            'Name,City|Alice,NYC',
        )
        assert code == 0
        e = WordEngine()
        e.open(str(out))
        assert e.extract_tables() == [[['Name', 'City'], ['Alice', 'NYC']]]

    def test_add_table_from_csv(self, tmp_path: Path) -> None:
        csv_file = tmp_path / 'data.csv'
        csv_file.write_text('a,b\n1,2\n', encoding='utf-8')
        out = tmp_path / 't.docx'
        code, _ = _run(
            '--force',
            '--create',
            '--word',
            '--output',
            str(out),
            '--add-table',
            f'@{csv_file}',
        )
        assert code == 0
        e = WordEngine()
        e.open(str(out))
        assert e.extract_tables() == [[['a', 'b'], ['1', '2']]]


class TestReverseCli:
    def test_markdown_to_word(self, tmp_path: Path) -> None:
        md = tmp_path / 'doc.md'
        md.write_text('# Hello\n\nSome **bold** text.\n', encoding='utf-8')
        out = tmp_path / 'doc.docx'
        code, _ = _run('--force', '--output', str(out), str(md))
        assert code == 0
        e = WordEngine()
        e.open(str(out))
        assert 'Hello' in e.extract_text()

    def test_html_to_word(self, tmp_path: Path) -> None:
        html = tmp_path / 'doc.html'
        html.write_text('<html><body><h1>Hi</h1><p>Some text.</p></body></html>', encoding='utf-8')
        out = tmp_path / 'doc.docx'
        code, _ = _run('--force', '--output', str(out), str(html))
        assert code == 0
        e = WordEngine()
        e.open(str(out))
        assert 'Hi' in e.extract_text()

    def test_json_to_excel(self, tmp_path: Path) -> None:
        data = [{'name': 'Alice', 'age': 30}, {'name': 'Bob', 'age': 25}]
        json_path = tmp_path / 'data.json'
        json_path.write_text(json.dumps(data), encoding='utf-8')
        out = tmp_path / 'data.xlsx'
        code, _ = _run('--force', '--output', str(out), str(json_path))
        assert code == 0
        from openpyxl import load_workbook

        ws = load_workbook(out).active
        rows = [list(r) for r in ws.iter_rows(values_only=True)]
        assert rows == [['name', 'age'], ['Alice', 30], ['Bob', 25]]


class TestBatchCli:
    def test_batch_replace(self, tmp_path: Path) -> None:
        for i in range(3):
            e = WordEngine()
            e.create()
            e.add_text(f'Report {i} final')
            e.save(str(tmp_path / f'doc{i}.docx'))
        code, out = _run(
            '--files',
            str(tmp_path / '*.docx'),
            '--force',
            '--replace',
            'final',
            '--replace-new',
            'v2',
        )
        assert code == 0
        assert '3 succeeded' in out
        e = WordEngine()
        e.open(str(tmp_path / 'doc0-out.docx'))
        assert 'v2' in e.extract_text()

    def test_batch_stdout_rejected(self, tmp_path: Path) -> None:
        e = WordEngine()
        e.create()
        e.save(str(tmp_path / 'doc0.docx'))
        code, out = _run('--files', str(tmp_path / '*.docx'), '--stdout')
        assert code == 2
        assert '--stdout' in out


class TestCompressMediaCli:
    def test_compress_media(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        img = tmp_path / 'big.jpg'
        Image.new('RGB', (3000, 1500), (90, 90, 90)).save(img, 'JPEG', quality=95)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(4))
        src = tmp_path / 'deck.pptx'
        prs.save(str(src))
        out = tmp_path / 'deck-out.pptx'
        code, out_text = _run(
            '--force',
            '--compress-media',
            '1600,60',
            '--output',
            str(out),
            str(src),
        )
        assert code == 0
        assert 'saved' in out_text.lower()
        assert out.exists()
        assert out.stat().st_size < src.stat().st_size


class TestOneShotOrdering:
    """File-first ordering regression (was broken under the Typer group)."""

    def test_file_before_options(self, tmp_path: Path) -> None:
        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.add_text('hello world')
        e.save(str(docx))
        code, out = _run(str(docx), '--replace', 'world', '--replace-new', 'there')
        assert code == 0
        assert 'Replaced 1 occurrence' in out
        e2 = WordEngine()
        e2.open(str(tmp_path / 't-out.docx'))
        assert 'hello there' in e2.extract_text()

    def test_options_before_file_still_works(self, tmp_path: Path) -> None:
        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.add_text('hello world')
        e.save(str(docx))
        code, out = _run('--replace', 'world', '--replace-new', 'there', str(docx))
        assert code == 0
        assert 'Replaced 1 occurrence' in out

    def test_file_first_with_add_table(self, tmp_path: Path) -> None:
        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.save(str(docx))
        code, _ = _run(str(docx), '--add-table', 'A,B|1,2')
        assert code == 0
        e2 = WordEngine()
        e2.open(str(tmp_path / 't-out.docx'))
        assert e2.extract_tables() == [[['A', 'B'], ['1', '2']]]


class TestOpenSubcommand:
    def test_open_dispatch(self, tmp_path: Path, monkeypatch) -> None:
        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.add_text('base')
        e.save(str(docx))
        calls = []
        monkeypatch.setattr('src.cli.repl.Prompt.ask', lambda *a, **k: calls.append(a) or 'quit')
        result = runner.invoke(open_app, [str(docx)])
        assert result.exit_code == 0
        assert 'Opened' in result.stdout
        assert calls  # the prompt was reached

    def test_open_main_cli_routes_to_open(self, tmp_path: Path, monkeypatch) -> None:
        import pytest

        import src.cli.main as cli

        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.add_text('base')
        e.save(str(docx))
        monkeypatch.setattr('src.cli.repl.Prompt.ask', lambda *a, **k: 'quit')
        monkeypatch.setattr('sys.argv', ['tianshang-scribe', 'open', str(docx)])
        with pytest.raises(SystemExit):
            cli.main_cli()

    def test_open_main_cli_routes_to_oneshot(self, tmp_path: Path, monkeypatch) -> None:
        import pytest

        import src.cli.main as cli

        docx = tmp_path / 't.docx'
        e = WordEngine()
        e.create()
        e.add_text('base')
        e.save(str(docx))
        monkeypatch.setattr(
            'sys.argv',
            ['tianshang-scribe', str(docx), '--replace', 'base', '--replace-new', 'new'],
        )
        with pytest.raises(SystemExit):
            cli.main_cli()
        e2 = WordEngine()
        e2.open(str(tmp_path / 't-out.docx'))
        assert 'new' in e2.extract_text()
