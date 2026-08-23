"""extract_presentation_data — Read-only semantic inspection of a .pptx deck.

``outline`` returns per-slide layout/title/bullets/notes/transition;
``structure`` returns the shape-type distribution per slide and overall;
``notes`` returns the full speaker-notes text per slide; ``master_info``
returns the slide-master/layout inventory with placeholder types.
"""

from __future__ import annotations

from pathlib import Path
from typing import Annotated, Any, Literal

from pptx import Presentation
from pydantic import Field

from tianshang_scribe.mcp.errors import McpErrorCode, error_response, success_response

_MODES = ('outline', 'notes', 'master_info', 'structure')

_TRANSITION_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def _slide_title(slide: Any) -> str:
    title = slide.shapes.title
    return title.text if title is not None else ''


def _slide_bullets(slide: Any) -> list[str]:
    """Non-empty paragraph texts from every non-title text-bearing shape."""
    bullets: list[str] = []
    title = slide.shapes.title
    for shape in slide.shapes:
        if not shape.has_text_frame or shape == title:
            continue
        for para in shape.text_frame.paragraphs:
            text = ''.join(run.text for run in para.runs).strip()
            if text:
                bullets.append(text)
    return bullets


def _slide_transition(slide: Any) -> str | None:
    """Read back the transition type written by ``PptEngine.set_transition``."""
    el = slide.element.find(f'{{{_TRANSITION_NS}}}transition')
    if el is None or len(el) == 0:
        return None
    tag: str = el[0].tag.rsplit('}', 1)[-1]
    return tag


def _placeholder_types(shapes: Any) -> list[str]:
    """Human-readable placeholder type list (``sldNum``/``ftr``/``dt`` audit)."""
    kinds: list[str] = []
    for shape in shapes:
        fmt = getattr(shape, 'placeholder_format', None)
        if fmt is None:
            continue
        raw = getattr(fmt, 'type', None)
        name = getattr(raw, 'name', None) or str(raw)
        kinds.append(f'{name}#{fmt.idx}')
    return kinds


def _master_inventory(prs: Any) -> list[dict[str, Any]]:
    masters: list[dict[str, Any]] = []
    for mi, master in enumerate(prs.slide_masters):
        layouts = [
            {
                'name': layout.name,
                'shape_count': len(layout.shapes),
                'placeholders': _placeholder_types(layout.placeholders),
            }
            for layout in master.slide_layouts
        ]
        masters.append(
            {
                'index': mi,
                'placeholders': _placeholder_types(master.placeholders),
                'layout_count': len(layouts),
                'layouts': layouts,
            }
        )
    return masters


def _shape_counts(slide: Any) -> dict[str, int]:
    counts = {'text': 0, 'table': 0, 'chart': 0, 'picture': 0, 'media': 0, 'other': 0}
    for shape in slide.shapes:
        if getattr(shape, 'has_table', False):
            counts['table'] += 1
        elif getattr(shape, 'has_chart', False):
            counts['chart'] += 1
        elif shape.shape_type is not None and 'MEDIA' in str(shape.shape_type):
            counts['media'] += 1
        elif str(shape.shape_type) == 'PICTURE (13)':
            counts['picture'] += 1
        elif shape.has_text_frame:
            counts['text'] += 1
        else:
            counts['other'] += 1
    return counts


def extract_presentation_data(
    input_path: Annotated[str, Field(description='Path to the existing .pptx presentation.')],
    mode: Annotated[
        Literal['outline', 'notes', 'master_info', 'structure'],
        Field(
            description=(
                "'outline' (per-slide semantics), 'structure' (shape-type census), "
                "'notes' (full speaker-notes text) or 'master_info' "
                '(masters/layouts inventory).'
            )
        ),
    ] = 'outline',
) -> dict[str, Any]:
    """Inspect a PowerPoint deck's structure without modifying it.

    Read-only — never modifies the input file. ``notes`` returns the full
    speaker-notes text per slide and ``master_info`` lists slide masters and
    layouts (useful to audit master-level placeholders). For workbook analysis
    use analyze_excel_data; to change slides use edit_presentation.
    """
    if not Path(input_path).exists():
        return error_response(McpErrorCode.DOCUMENT_NOT_FOUND, f"'{input_path}' not found.")
    if not input_path.lower().endswith('.pptx'):
        return error_response(
            McpErrorCode.UNSUPPORTED_FORMAT,
            f"'{input_path}' is not a .pptx presentation.",
        )

    try:
        prs = Presentation(input_path)
        payload: dict[str, Any] = {'input_path': input_path, 'mode': mode}
        if mode == 'master_info':
            payload['slide_count'] = len(prs.slides)
            payload['masters'] = _master_inventory(prs)
            return success_response(payload)

        slides: list[dict[str, Any]] = []
        if mode == 'outline':
            for idx, slide in enumerate(prs.slides):
                slides.append(
                    {
                        'index': idx,
                        'layout': slide.slide_layout.name,
                        'title': _slide_title(slide),
                        'bullets': _slide_bullets(slide),
                        'notes': (
                            slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ''
                        ),
                        'transition': _slide_transition(slide),
                    }
                )
        elif mode == 'notes':
            for idx, slide in enumerate(prs.slides):
                notes_text = (
                    slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ''
                )
                slides.append(
                    {
                        'index': idx,
                        'title': _slide_title(slide),
                        'has_notes': bool(notes_text),
                        'notes': notes_text,
                    }
                )
        else:
            totals = {'text': 0, 'table': 0, 'chart': 0, 'picture': 0, 'media': 0, 'other': 0}
            for idx, slide in enumerate(prs.slides):
                counts = _shape_counts(slide)
                for key in totals:
                    totals[key] += counts[key]
                slides.append({'index': idx, **counts})
            payload_extra: dict[str, Any] = {'totals': totals, 'slide_count': len(prs.slides)}
        payload['slides'] = slides
        if mode == 'structure':
            payload.update(payload_extra)
        return success_response(payload)
    except Exception as e:  # surface any read failure as a structured error
        return error_response(McpErrorCode.INTERNAL_ERROR, str(e))
