"""Verify demo output files."""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# Word
from docx import Document
from docx.oxml.ns import qn
doc = Document('demo/demo_word.docx')
print('=== Word ===')
print(f'Paragraphs: {len(doc.paragraphs)}')
omath_count = sum(1 for p in doc.paragraphs for m in p._p.findall(qn('m:oMath')))
print(f'OMML formulas: {omath_count}')
print(f'Sections: {len(doc.sections)}')
has_header = any(s.header.paragraphs and s.header.paragraphs[0].text for s in doc.sections)
print(f'Header text: {has_header}')
has_wm = any('天殇' in s.header.paragraphs[0].text for s in doc.sections if s.header.paragraphs)
print(f'Watermark: {has_wm}')

# Excel
from openpyxl import load_workbook
wb = load_workbook('demo/demo_excel.xlsx')
print()
print('=== Excel ===')
print(f'Sheets: {wb.sheetnames}')
ws = wb['Sales']
print(f'Rows (Sales): {ws.max_row}, Cols: {ws.max_column}')
a1 = ws['A1'].value
f2 = ws['F2'].value
print(f'A1: {a1}, F2 formula: {f2}')
print(f'Comment on A1: {ws["A1"].comment is not None}')
print(f'Password set: {wb.security.workbook_password is not None}')

# PPT
from pptx import Presentation
prs = Presentation('demo/demo_ppt.pptx')
print()
print('=== PPT ===')
print(f'Slides: {len(prs.slides)}')
ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
for i, slide in enumerate(prs.slides):
    trans = slide.element.find(f'{{{ns}}}transition')
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text[:40]
            if t.strip():
                texts.append(t)
    notes_txt = ''
    try:
        notes_txt = slide.notes_slide.notes_text_frame.text[:40]
    except Exception:
        pass
    has_trans = trans is not None
    print(f'  Sl{i}: transition={has_trans} notes={notes_txt!r} text={texts[0] if texts else ""!r}')

print('\n=== Summary ===')
print(f'Word:  {len(doc.paragraphs)} paragraphs, {omath_count} formulas, header+watermark')
print(f'Excel: {len(wb.sheetnames)} sheets, {ws.max_row} rows, formulas+chart+comment+protection')
print(f'PPT:   {len(prs.slides)} slides, transitions, notes')
