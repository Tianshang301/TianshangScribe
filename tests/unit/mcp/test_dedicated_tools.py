"""Baseline tests for the dedicated MCP tools (v0.8.0 expansion)."""

from __future__ import annotations

from pathlib import Path

from tianshang_scribe.core.document import DocumentType, create_document
from tianshang_scribe.mcp.tools.analyze_excel import analyze_excel_data


def _make_workbook(path: Path) -> None:
    e = create_document(DocumentType.EXCEL)
    e.add_sheet('Data')
    ws = e.wb['Data']
    ws.append(['name', 'score', 'group'])
    ws.append(['alice', 10, 'a'])
    ws.append(['bob', 20, 'b'])
    ws.append(['alice', 10, 'a'])  # duplicate row
    ws.append(['carol', 30, 'a'])
    e.save(str(path))


def test_analyze_excel_data_baseline(tmp_path: Path) -> None:
    book = tmp_path / 'data.xlsx'
    _make_workbook(book)
    res = analyze_excel_data(str(book))
    assert res['success'] is True, res
    data = res['data']
    assert data['sheet_count'] == 2  # default 'Sheet' + 'Data'
    assert data['duplicate_row_count'] == 1
    sheet = next(s for s in data['sheets'] if s['name'] == 'Data')
    assert sheet['row_count'] == 4
    assert sheet['headers'] == ['name', 'score', 'group']
    by_name = {c['name']: c for c in sheet['columns']}
    assert by_name['score']['inferred_type'] == 'numeric'
    assert by_name['score']['numeric']['min'] == 10
    assert by_name['score']['numeric']['max'] == 30
    assert by_name['group']['inferred_type'] == 'categorical'
    assert by_name['name']['null_count'] == 0


def test_analyze_excel_rejects_non_excel(tmp_path: Path) -> None:
    from tianshang_scribe.core.document import DocumentType
    from tianshang_scribe.mcp.tools.analyze_excel import analyze_excel_data

    doc = tmp_path / 'a.docx'
    e = create_document(DocumentType.WORD)
    e.add_text('hello')
    e.save(str(doc))
    res = analyze_excel_data(str(doc))
    assert res['success'] is False
    assert res['error_code'] == 1003  # UNSUPPORTED_FORMAT


def test_create_excel_workbook_baseline(tmp_path: Path) -> None:
    from openpyxl import load_workbook

    from tianshang_scribe.mcp.tools.excel_create import create_excel_workbook

    out = tmp_path / 'book.xlsx'
    res = create_excel_workbook(
        str(out),
        sheets=[
            {
                'name': 'Data',
                'headers': ['name', 'score'],
                'rows': [['alice', 10], ['bob', 20]],
                'formulas': {'C1': '=SUM(B2:B3)'},
                'number_format': 'B2:B3=0.00',
            }
        ],
    )
    assert res['success'] is True, res
    assert res['data']['sheet_count'] == 1
    wb = load_workbook(str(out))
    assert 'Data' in wb.sheetnames
    ws = wb['Data']
    assert ws['A1'].value == 'name'
    assert ws['C1'].value == '=SUM(B2:B3)'
    assert ws['B2'].number_format == '0.00'


def test_edit_excel_workbook_write_cell_and_sheet(tmp_path: Path) -> None:
    from openpyxl import load_workbook

    from tianshang_scribe.mcp.tools.excel_create import create_excel_workbook
    from tianshang_scribe.mcp.tools.excel_edit import edit_excel_workbook

    out = tmp_path / 'book.xlsx'
    create_excel_workbook(str(out), sheets=[{'name': 'Data', 'rows': [['alice', 10]]}])
    res = edit_excel_workbook(
        str(out),
        operations=[
            {'action': 'write_cell', 'cell': 'C1', 'value': 42, 'sheet_name': 'Data'},
            {'action': 'add_sheet', 'sheet_name': 'Extra'},
        ],
    )
    assert res['success'] is True, res
    wb = load_workbook(str(out))
    assert wb['Data']['C1'].value == 42
    assert 'Extra' in wb.sheetnames


def test_create_presentation_baseline(tmp_path: Path) -> None:
    from pptx import Presentation

    from tianshang_scribe.mcp.tools.ppt_create import create_presentation

    out = tmp_path / 'deck.pptx'
    res = create_presentation(
        str(out),
        slides=[
            {'title': 'Intro', 'bullets': ['a', 'b'], 'layout': 'Title and Content'},
        ],
    )
    assert res['success'] is True, res
    assert res['data']['slide_count'] == 1
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_edit_presentation_add_slide(tmp_path: Path) -> None:
    from pptx import Presentation

    from tianshang_scribe.mcp.tools.ppt_create import create_presentation
    from tianshang_scribe.mcp.tools.ppt_edit import edit_presentation

    out = tmp_path / 'deck.pptx'
    create_presentation(str(out), slides=[{'title': 'First'}])
    res = edit_presentation(
        str(out),
        operations=[{'action': 'add_slide'}, {'action': 'add_text', 'text': 'hello', 'slide_index': 1}],
    )
    assert res['success'] is True, res
    prs = Presentation(str(out))
    assert len(prs.slides) == 2

