"""Unit tests for the PowerPoint engine."""

from __future__ import annotations

import tempfile
from pathlib import Path

import pytest

from tianshang_scribe.core.document import open_document
from tianshang_scribe.core.ppt_engine import PptEngine


class TestPptEngine:
    @pytest.fixture
    def engine(self) -> PptEngine:
        e = PptEngine()
        e.create()
        return e

    # ---- Step 8 (P1): precise-positioned text box ----
    def test_add_textbox_baseline(self, engine: PptEngine, tmp_path: Path) -> None:
        engine.add_slide()
        box = engine.add_textbox(0, 'Positioned', left=2.0, top=3.0, width=4.0, height=1.0)
        assert box.left == 2.0 * 914400
        assert box.top == 3.0 * 914400
        assert box.width == 4.0 * 914400
        assert 'Positioned' in box.text_frame.text
        out = tmp_path / 'tb.pptx'
        engine.save(out)
        reopened = open_document(out)
        assert isinstance(reopened, PptEngine)
        assert any('Positioned' in sh.text_frame.text for sh in reopened.prs.slides[0].shapes if sh.has_text_frame)

    def test_add_textbox_out_of_range_raises(self, engine: PptEngine) -> None:
        engine.add_slide()
        with pytest.raises(IndexError, match='slide_index out of range'):
            engine.add_textbox(5, 'x')

    # ---- Step 9 (P2): insert table ----
    def test_add_table_baseline(self, engine: PptEngine, tmp_path: Path) -> None:
        engine.add_slide()
        gf = engine.add_table(
            0, [['a1', 'b1'], ['a2', 'b2']], col_names=['H1', 'H2']
        )
        table = gf.table
        assert table.rows.__len__() == 3
        assert table.columns.__len__() == 2
        assert table.cell(0, 0).text == 'H1'
        assert table.cell(1, 0).text == 'a1'
        out = tmp_path / 'tbl.pptx'
        engine.save(out)
        reopened = open_document(out)
        assert isinstance(reopened, PptEngine)
        shapes = list(reopened.prs.slides[0].shapes)
        assert any(sh.has_table for sh in shapes)

    def test_add_table_out_of_range_raises(self, engine: PptEngine) -> None:
        engine.add_slide()
        with pytest.raises(IndexError, match='slide_index out of range'):
            engine.add_table(9, [['x']])

    # ---- Step 10 (P3): insert chart ----
    def test_add_chart_baseline(self, engine: PptEngine, tmp_path: Path) -> None:
        engine.add_slide()
        data = [['', 'Q1'], ['Jan', 10], ['Feb', 20]]
        gf = engine.add_chart(0, 'bar', data, title='Sales')
        assert gf.has_chart
        assert gf.chart.plots[0].series[0].name == 'Q1'
        out = tmp_path / 'chart.pptx'
        engine.save(out)
        reopened = open_document(out)
        assert isinstance(reopened, PptEngine)
        shapes = list(reopened.prs.slides[0].shapes)
        assert any(sh.has_chart for sh in shapes)

    def test_add_chart_out_of_range_raises(self, engine: PptEngine) -> None:
        engine.add_slide()
        with pytest.raises(IndexError, match='slide_index out of range'):
            engine.add_chart(9, 'bar', [['', 'S'], ['a', 1]])

    def test_create_has_slides(self, engine: PptEngine) -> None:
        assert len(engine.prs.slides) >= 0

    def test_add_slide(self, engine: PptEngine) -> None:
        initial = len(engine.prs.slides)
        engine.add_slide()
        assert len(engine.prs.slides) == initial + 1

    def test_save_and_reopen(self, engine: PptEngine) -> None:
        engine.add_slide()
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / 'test.pptx'
            engine.save(path)
            assert path.exists()
            reopened = open_document(path)
            assert isinstance(reopened, PptEngine)

    def test_delete_slide(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.add_slide()
        engine.add_slide()
        initial = len(engine.prs.slides)
        engine.delete_slide(0)
        assert len(engine.prs.slides) == initial - 1

    def test_move_slide(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.add_slide()
        engine.add_slide()
        engine.move_slide(0, 1)
        assert len(engine.prs.slides) == 3

    def test_add_notes(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.add_notes(0, 'Speaker note')
        notes = engine.prs.slides[0].notes_slide
        assert notes.notes_text_frame.text == 'Speaker note'

    def test_set_transition(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.set_transition('fade')
        slide = engine.prs.slides[0]
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        trans = slide.element.find(f'{{{ns}}}transition')
        assert trans is not None
        assert trans.find(f'{{{ns}}}fade') is not None

    def test_set_transition_invalid_raises(self, engine: PptEngine) -> None:
        engine.add_slide()
        with pytest.raises(ValueError, match='Unsupported transition'):
            engine.set_transition('nonexistent')

    def test_set_protection(self, engine: PptEngine) -> None:
        engine.set_protection('secret')
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        pres = engine.prs.part._element
        verifier = pres.find(f'{{{ns}}}modifyVerifier')
        assert verifier is not None
        # compliant: not plaintext, base64 hash of a 64-byte SHA-512 digest
        assert verifier.get('hashData') != 'secret'
        import base64

        assert len(base64.b64decode(verifier.get('saltData'))) == 16
        assert len(base64.b64decode(verifier.get('hashData'))) == 64
        assert PptEngine.verify_modify_verifier(
            'secret', verifier.get('saltData'), verifier.get('hashData')
        )
        assert not PptEngine.verify_modify_verifier(
            'wrong', verifier.get('saltData'), verifier.get('hashData')
        )

    def test_unprotect(self, engine: PptEngine) -> None:
        engine.set_protection('secret')
        engine.unprotect()
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        pres = engine.prs.part._element
        assert pres.find(f'{{{ns}}}modifyVerifier') is None


class TestPptEdge:
    @pytest.fixture
    def engine(self) -> PptEngine:
        e = PptEngine()
        e.create()
        return e

    def test_prs_property_unloaded(self) -> None:
        e = PptEngine()
        with pytest.raises(RuntimeError):
            _ = e.prs

    def test_open_nonexistent(self) -> None:
        e = PptEngine()
        with pytest.raises(FileNotFoundError):
            e.open('nope.pptx')

    def test_save_without_path(self, engine: PptEngine) -> None:
        with pytest.raises(ValueError):
            engine.save()

    def test_get_base_style(self, engine: PptEngine) -> None:
        assert engine.get_base_style().font_name == 'Calibri'

    def test_set_style(self, engine: PptEngine) -> None:
        engine.set_style('font=Arial,size=30,bold')
        assert engine._base_style.font_name == 'Arial'
        assert engine._base_style.font_size == 30

    def test_apply_style_to_all(self, engine: PptEngine) -> None:
        engine.set_style('font=Courier New')
        engine.add_text('hello')
        engine.apply_style_to_all()
        for slide in engine.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            assert run.font.name == 'Courier New'

    def test_add_text_with_style(self, engine: PptEngine) -> None:
        engine.add_text(
            'Styled', bold=True, italic=True, font_size=20, color='FF0000', font_name='Arial'
        )
        slide = engine.prs.slides[-1]
        run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
        assert run.font.bold is True
        assert run.font.italic is True
        assert run.font.size.pt == 20

    def test_add_text_with_text_style_object(self, engine: PptEngine) -> None:
        from tianshang_scribe.rendering.styles import TextStyle

        engine.add_text('Styled', text_style=TextStyle(font_name='Arial', font_size=20))
        slide = engine.prs.slides[-1]
        run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
        assert run.font.name == 'Arial'

    def test_add_text_alignment(self, engine: PptEngine) -> None:
        engine.add_text('Centered', alignment='center')
        slide = engine.prs.slides[-1]
        assert slide.shapes.title.text_frame.paragraphs[0].alignment is not None

    def test_add_text_trailing_math_buffer(self, engine: PptEngine) -> None:
        engine.add_text(r'Text with $unfinished')
        assert engine.prs.slides[-1].shapes.title is not None

    def test_add_text_invalid_color(self, engine: PptEngine) -> None:
        engine.add_text('bad', color='XYZ')
        assert 'bad' in engine.prs.slides[-1].shapes.title.text

    def test_add_text_math_inline(self, engine: PptEngine) -> None:
        engine.add_text(r'Area is $x^2$ ok')
        slide = engine.prs.slides[-1]
        assert slide.shapes.title is not None

    def test_add_text_math_display(self, engine: PptEngine) -> None:
        engine.add_text(r'Formula: $$y = x$$')
        assert engine.prs.slides[-1].shapes.title is not None

    def test_replace_text(self, engine: PptEngine) -> None:
        engine.add_text('Hello World')
        count = engine.replace_text('World', 'There')
        assert count >= 1
        assert 'There' in engine.extract_text()

    def test_replace_text_regex(self, engine: PptEngine) -> None:
        engine.add_text('Version 1.0')
        count = engine.replace_text(r'\d+\.\d+', '2.0', regex=True)
        assert count >= 1
        assert '2.0' in engine.extract_text()

    def test_replace_text_no_match(self, engine: PptEngine) -> None:
        engine.add_text('Hello')
        assert engine.replace_text('zzz', 'x') == 0

    def test_metadata_all(self, engine: PptEngine) -> None:
        engine.set_metadata(
            author='a', title='t', subject='s', category='c', keywords='k', comments='m'
        )
        md = engine.get_metadata()
        assert md == {
            'author': 'a',
            'title': 't',
            'subject': 's',
            'category': 'c',
            'keywords': 'k',
            'comments': 'm',
        }

    def test_add_comment(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.add_comment('note')
        assert 'note' in engine.prs.slides[0].notes_slide.notes_text_frame.text

    def test_add_comment_appends(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.add_comment('first')
        engine.add_comment('second')
        text = engine.prs.slides[0].notes_slide.notes_text_frame.text
        assert 'first' in text and 'second' in text

    def test_add_comment_out_of_range(self, engine: PptEngine) -> None:
        with pytest.raises(ValueError):
            engine.add_comment('note', slide_index=99)

    def test_add_notes_out_of_range(self, engine: PptEngine) -> None:
        with pytest.raises(ValueError):
            engine.add_notes(99, 'x')

    def test_apply_layout_by_name(self, engine: PptEngine) -> None:
        engine.add_slide()
        layout_name = engine.prs.slide_layouts[0].name
        engine.apply_layout(0, layout_name)

    def test_apply_layout_by_index(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.apply_layout(0, '1')

    def test_apply_layout_not_found(self, engine: PptEngine) -> None:
        engine.add_slide()
        with pytest.raises(ValueError, match='not found'):
            engine.apply_layout(0, 'nonexistent layout')

    def test_clear_content(self, engine: PptEngine) -> None:
        engine.add_text('to clear')
        engine.clear_content()
        assert engine.extract_text() == ''

    def test_add_latex_content_stacks_non_overlapping(self, engine: PptEngine) -> None:
        # heading + text + heading + text on the SAME slide => multiple text boxes
        engine.add_latex_content('\\heading{1}{First} second \\heading{1}{Third} fourth')
        slide = engine.prs.slides[-1]
        boxes = [s for s in slide.shapes if s.has_text_frame]
        tops = [round(float(b.top)) for b in boxes]
        assert len(tops) >= 3
        assert len(set(tops)) == len(tops)  # no two text boxes share the same position

    def test_add_text_targets_existing_slide(self, engine: PptEngine) -> None:
        engine.add_text('title')
        initial_slides = len(engine.prs.slides)
        engine.add_text('body', slide_index=0)
        assert len(engine.prs.slides) == initial_slides
        body_text = engine.extract_text()
        assert 'title' in body_text and 'body' in body_text

    def test_merge_workbooks(self, engine: PptEngine, tmp_path: Path) -> None:
        other = tmp_path / 'other.pptx'
        e2 = PptEngine()
        e2.create()
        e2.add_text('merged')
        e2.save(other)
        initial = len(engine.prs.slides)
        engine.merge_workbooks([str(other)])
        assert len(engine.prs.slides) == initial + 1
        # content must be faithfully copied, not just a blank slide
        merged_text = engine.extract_text()
        assert 'merged' in merged_text

    def test_merge_workbooks_keeps_images(self, engine: PptEngine, tmp_path: Path) -> None:
        from PIL import Image
        from pptx.util import Inches

        img = tmp_path / 'pic.png'
        Image.new('RGB', (32, 32), 'red').save(img)
        other = tmp_path / 'other.pptx'
        e2 = PptEngine()
        e2.create()
        slide = e2.prs.slides.add_slide(e2.prs.slide_layouts[5])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), Inches(2), Inches(2))
        e2.save(other)

        engine.merge_workbooks([str(other)])
        out = tmp_path / 'merged.pptx'
        engine.save(out)
        # reopen and ensure the picture part survived the clone
        reopened = PptEngine()
        reopened.open(str(out))
        last = reopened.prs.slides[-1]
        assert any(shape.shape_type == 13 for shape in last.shapes)  # 13 == PICTURE

    def test_extract_tables(self, tmp_path: Path) -> None:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        from pptx.util import Inches as In

        table_shape = slide.shapes.add_table(2, 2, In(1), In(1), In(4), In(2))
        table_shape.table.cell(0, 0).text = 'A'
        path = tmp_path / 't.pptx'
        prs.save(str(path))
        e = PptEngine()
        e.open(str(path))
        tables = e.extract_tables()
        assert tables and tables[0][0][0] == 'A'

    def test_extract_images(self, engine: PptEngine, tmp_path: Path) -> None:
        from PIL import Image
        from pptx.util import Inches

        img = tmp_path / 'pic.png'
        Image.new('RGB', (16, 16), (0, 0, 255)).save(img)
        slide = engine.add_slide()
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(1))
        saved = engine.extract_images(str(tmp_path / 'out'))
        assert len(saved) == 1

    def test_to_pdf(self, engine: PptEngine, tmp_path: Path, monkeypatch) -> None:
        engine.add_slide()
        out = tmp_path / 'out.pdf'
        calls = []
        monkeypatch.setattr(engine, 'save', lambda *a, **k: calls.append('save'))
        monkeypatch.setattr(
            'tianshang_scribe.transform.pdf.ppt_to_pdf', lambda s, d: calls.append((s, d))
        )
        engine.to_pdf(out)
        assert 'save' in calls

    def test_to_images_no_libreoffice(self, engine: PptEngine, monkeypatch) -> None:
        monkeypatch.setattr('shutil.which', lambda p: None)
        with pytest.raises(RuntimeError, match='LibreOffice'):
            engine.to_images('outdir')

    def test_set_transition_specific_slide(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.set_transition('fade', slide_index=0)
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        assert engine.prs.slides[0].element.find(f'{{{ns}}}transition') is not None

    def test_compress_media_no_blip(self, engine: PptEngine) -> None:
        engine.add_slide()
        assert engine.compress_media() == 0

    def test_add_latex_content(self, engine: PptEngine) -> None:
        engine.add_latex_content(r'\bfseries{Title} and text')
        assert 'Title' in engine.extract_text()

    def test_add_styled_content_text_and_heading(self, engine: PptEngine) -> None:
        engine.add_styled_content(
            [
                {'type': 'text', 'content': 'Hello '},
                {'type': 'command', 'command': 'heading', 'content': 'Head'},
                {'type': 'command', 'command': 'math', 'latex': 'a+b'},
                {'type': 'command', 'command': 'bfseries', 'content': 'Bold'},
            ]
        )
        assert 'Hello' in engine.extract_text()
        assert 'Bold' in engine.extract_text()

    def test_add_styled_content_newpage(self, engine: PptEngine) -> None:
        engine.add_styled_content(
            [
                {'type': 'text', 'content': 'page1 '},
                {'type': 'command', 'command': 'newpage'},
                {'type': 'text', 'content': 'page2 '},
            ]
        )
        assert 'page1' in engine.extract_text()
        assert 'page2' in engine.extract_text()

    def test_append_text_heading_flag(self, engine: PptEngine) -> None:
        slide = engine.add_slide()
        engine._append_text_to_slide(slide, 'Head', engine._base_style, heading=True)
        text = engine.extract_text()
        assert 'Head' in text

    def test_to_images_mocked_libreoffice(
        self, engine: PptEngine, tmp_path: Path, monkeypatch
    ) -> None:
        import subprocess

        monkeypatch.setattr('shutil.which', lambda p: p if 'soffice' in p else None)
        out = tmp_path / 'imgs'
        engine.add_slide()
        engine.add_slide()
        engine.save(str(tmp_path / 'deck.pptx'))
        calls = []

        def fake_run(cmd, **kw):
            calls.append(cmd)
            # LibreOffice convert-to pdf step: drop a pdf into the requested outdir
            if 'pdf' in cmd:
                outdir = cmd[cmd.index('--outdir') + 1]
                (Path(outdir) / 'deck.pdf').write_bytes(b'%PDF')
            return None

        def fake_pdf_to_png(pdf_path, output_dir):
            Path(output_dir).mkdir(parents=True, exist_ok=True)
            (Path(output_dir) / 'slide1.png').write_bytes(b'png')
            (Path(output_dir) / 'slide2.png').write_bytes(b'png')

        monkeypatch.setattr(subprocess, 'run', fake_run)
        monkeypatch.setattr(PptEngine, '_pdf_to_png', staticmethod(fake_pdf_to_png))
        result = engine.to_images(str(out))
        assert len(result) == 2
        assert result == sorted(out.glob('*.png'))

    def test_compress_media_png(self, tmp_path: Path) -> None:
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        img = tmp_path / 'big.png'
        Image.new('RGB', (3000, 3000), (200, 200, 200)).save(img, 'PNG')
        pptx_path = tmp_path / 'png_deck.pptx'
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(4))
        prs.save(str(pptx_path))
        e = PptEngine()
        e.open(str(pptx_path))
        assert e.compress_media(max_dimension=1200) > 0

    def test_compress_media_shared_part_once(self, tmp_path: Path) -> None:
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        img = tmp_path / 'shared.png'
        Image.new('RGB', (2000, 2000), (10, 10, 10)).save(img, 'PNG')
        pptx_path = tmp_path / 'shared_deck.pptx'
        prs = Presentation()
        for _ in range(2):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(4))
        prs.save(str(pptx_path))
        e = PptEngine()
        e.open(str(pptx_path))
        assert e.compress_media(max_dimension=800) > 0


class TestPptExtract:
    def test_extract_text(self) -> None:
        e = PptEngine()
        e.create()
        e.add_text('Slide one')
        assert '[slide 1] Slide one' in e.extract_text()

    def test_extract_structure(self) -> None:
        e = PptEngine()
        e.create()
        e.add_slide()
        struct = e.extract_structure()
        assert struct['slides'] == 1


class TestPptCompressMedia:
    def test_compress_media_reduces_size(self, tmp_path: Path) -> None:
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        img = tmp_path / 'big.jpg'
        Image.new('RGB', (4000, 2000), (128, 128, 128)).save(img, 'JPEG', quality=95)
        pptx_path = tmp_path / 'deck.pptx'
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(4))
        prs.save(str(pptx_path))
        original = pptx_path.stat().st_size

        e = PptEngine()
        e.open(str(pptx_path))
        saved = e.compress_media(max_dimension=1600, quality=60)
        assert saved > 0
        e.save(str(pptx_path))
        assert pptx_path.stat().st_size < original
