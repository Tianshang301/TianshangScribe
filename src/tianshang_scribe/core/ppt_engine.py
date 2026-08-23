"""PowerPoint presentation engine built on python-pptx."""

from __future__ import annotations

import base64
import copy
import hashlib
import io
import os
import re
import struct
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as _PresentationType
from pptx.util import Pt

from tianshang_scribe.core.document import DocumentABC
from tianshang_scribe.rendering.styles import TextStyle

#: Extension → MIME type for media inserted via ``add_movie`` / ``add_audio``.
_MEDIA_MIME: dict[str, str] = {
    '.mp4': 'video/mp4',
    '.mov': 'video/quicktime',
    '.avi': 'video/x-msvideo',
    '.mkv': 'video/x-matroska',
    '.mp3': 'audio/mpeg',
    '.wav': 'audio/x-wav',
    '.m4a': 'audio/mp4',
}

#: Master placeholder type → (standard idx, display name).
_MASTER_PH_SPECS: dict[str, tuple[str, str]] = {
    'sldNum': ('12', 'Slide Number Placeholder'),
    'ftr': ('11', 'Footer Placeholder'),
    'dt': ('10', 'Date Placeholder'),
}

#: Built-in themes for ``apply_theme``: slot → srgb hex (dk1/lt1 follow OOXML
#: semantics — dk1 is text-on-background, lt1 the background itself).
_THEME_PALETTES: dict[str, dict[str, str]] = {
    'office': {
        'dk1': '000000',
        'lt1': 'FFFFFF',
        'dk2': '44546A',
        'lt2': 'E7E6E6',
        'accent1': '4472C4',
        'accent2': 'ED7D31',
        'accent3': 'A5A5A5',
        'accent4': 'FFC000',
        'accent5': '5B9BD5',
        'accent6': '70AD47',
        'hlink': '0563C1',
        'folHlink': '954F72',
        'major_latin': 'Calibri Light',
        'minor_latin': 'Calibri',
    },
    'dark': {
        'dk1': 'FFFFFF',
        'lt1': '202124',
        'dk2': 'C9CDD6',
        'lt2': '2B2B3A',
        'accent1': '7AA2F7',
        'accent2': 'BB9AF7',
        'accent3': '9ECE6A',
        'accent4': 'FFA06A',
        'accent5': '7DCFFF',
        'accent6': 'E0AF68',
        'hlink': '7AA2F7',
        'folHlink': 'BB9AF7',
        'major_latin': 'Calibri Light',
        'minor_latin': 'Calibri',
    },
}

_DRAWINGML_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


class PptEngine(DocumentABC):
    """PowerPoint presentation engine: create, edit and style decks."""

    def __init__(self, path: str | Path | None = None) -> None:
        """Initialize the engine with an optional presentation path."""
        super().__init__(path)
        self._prs: _PresentationType | None = None
        self._base_style: TextStyle = TextStyle.default_ppt()
        self._text_cursors: dict[int, float] = {}

    @property
    def prs(self) -> _PresentationType:
        """Return the loaded presentation, raising if none is open."""
        if self._prs is None:
            raise RuntimeError('No presentation loaded. Call create() or open() first.')
        return self._prs

    def create(self) -> None:
        """Create a new blank presentation."""
        self._prs = Presentation()
        self._path = None
        self._base_style = TextStyle.default_ppt()

    def open(self, path: str | Path) -> None:
        """Open an existing presentation from the given path."""
        self._path = Path(path)
        if not self._path.exists():
            raise FileNotFoundError(f'File not found: {self._path}')
        self._prs = Presentation(str(self._path))
        self._base_style = TextStyle.default_ppt()

    def save(self, path: str | Path | None = None) -> None:
        """Save the presentation to the given path or the current one."""
        if path is not None:
            self._path = Path(path)
        if self._path is None:
            raise ValueError('No output path specified.')
        self._path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(self._path))

    def get_base_style(self) -> TextStyle:
        """Return the engine's base text style."""
        return self._base_style

    def add_text(
        self,
        text: str,
        bold: bool = False,
        italic: bool = False,
        font_name: str | None = None,
        font_size: int | None = None,
        color: str | None = None,
        alignment: str | None = None,
        text_style: TextStyle | None = None,
        slide_index: int | None = None,
        **kwargs: Any,
    ) -> Any:
        """Add text to a slide.

        By default a new slide is created and the text placed in its title
        placeholder. When ``slide_index`` is given, the text is appended to that
        existing slide (preferring its body placeholder, else a new text box).
        Supports inline math markup such as ``$E=mc^2$``.
        """
        inline = TextStyle(
            bold=bold or None,
            italic=italic or None,
            font_name=font_name,
            font_size=font_size,
            color=color,
            alignment=alignment,
        )
        final = self._base_style
        if text_style is not None:
            final = final.merge(text_style)
        final = final.merge(inline)

        if slide_index is not None and 0 <= slide_index < len(self.prs.slides):
            slide = self.prs.slides[slide_index]
            body = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape is not slide.shapes.title:
                    body = shape
                    break
            if body is not None and body.text_frame.text.strip():
                tf = self._place_textbox(slide).text_frame
            elif body is not None:
                tf = body.text_frame
            else:
                tf = self._place_textbox(slide).text_frame
            self._add_text_with_math(tf, text, final)
            return slide

        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title

        if not title_shape:
            return slide

        tf = title_shape.text_frame
        self._add_text_with_math(tf, text, final)
        return slide

    def add_textbox(
        self,
        slide_index: int,
        text: str,
        left: float = 1.0,
        top: float | None = None,
        width: float | None = None,
        height: float = 1.0,
        bold: bool = False,
        italic: bool = False,
        font_name: str | None = None,
        font_size: int | None = None,
        color: str | None = None,
        alignment: str | None = None,
        text_style: TextStyle | None = None,
    ) -> Any:
        """Add a text box at precise (inch) coordinates on the given slide.

        Unlike :meth:`add_text`, this places content at an explicit position
        rather than into the title/body placeholder. When ``top`` is omitted
        the box auto-stacks below previously placed boxes on the same slide
        (via the internal text cursor), so consecutive calls never overlap.
        """
        from pptx.util import Inches

        if not (0 <= slide_index < len(self.prs.slides)):
            raise IndexError(f'slide_index out of range: {slide_index}')
        slide = self.prs.slides[slide_index]
        if top is None:
            box = self._place_textbox(slide, left=left, width=width, height=height)
        else:
            if width is None:
                slide_width_in = (self.prs.slide_width or 914400 * 10) / 914400
                width = max(1.0, slide_width_in - 2 * left)
            box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        style = self._base_style
        if text_style is not None:
            style = style.merge(text_style)
        style = style.merge(
            TextStyle(
                bold=bold or None,
                italic=italic or None,
                font_name=font_name,
                font_size=font_size,
                color=color,
                alignment=alignment,
            )
        )
        self._add_text_with_math(box.text_frame, text, style)
        return box

    def add_table(
        self,
        slide_index: int,
        rows: list[list[Any]],
        col_names: list[Any] | None = None,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 8.0,
        height: float | None = None,
    ) -> Any:
        """Insert a table on ``slide_index``; ``col_names`` (if given) is the bold header row."""
        from pptx.util import Inches

        if not (0 <= slide_index < len(self.prs.slides)):
            raise IndexError(f'slide_index out of range: {slide_index}')
        if not rows and not col_names:
            raise ValueError('add_table requires rows and/or col_names')
        slide = self.prs.slides[slide_index]
        ncols = len(col_names) if col_names else len(rows[0])
        nrows = len(rows) + (1 if col_names else 0)
        if height is None:
            height = 0.4 * nrows
        gf = slide.shapes.add_table(
            nrows, ncols, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        table = gf.table
        r = 0
        if col_names:
            for c, name in enumerate(col_names):
                cell = table.cell(0, c)
                cell.text = str(name)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
            r = 1
        for row in rows:
            for c, val in enumerate(row):
                table.cell(r, c).text = str(val)
            r += 1
        return gf

    def add_chart(
        self,
        slide_index: int,
        chart_type: str,
        data: list[list[Any]],
        left: float = 1.0,
        top: float = 1.0,
        width: float = 6.0,
        height: float = 4.0,
        title: str | None = None,
    ) -> Any:
        """Insert a chart on ``slide_index``.

        ``data`` is a table where ``data[0]`` holds series names (``data[0][0]``
        is ignored) and each subsequent row is ``[category, *series_values]``.
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        if not (0 <= slide_index < len(self.prs.slides)):
            raise IndexError(f'slide_index out of range: {slide_index}')
        slide = self.prs.slides[slide_index]
        chart_data = CategoryChartData()
        chart_data.categories = [row[0] for row in data[1:]]
        for j in range(1, len(data[0])):
            name = data[0][j]
            values = [row[j] for row in data[1:]]
            chart_data.add_series(name, values)
        type_map = {
            'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'area': XL_CHART_TYPE.AREA,
            'doughnut': XL_CHART_TYPE.DOUGHNUT,
        }
        ct = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        gf = slide.shapes.add_chart(
            ct, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
        )
        if title:
            gf.chart.chart_title.text_frame.text = title
        return gf

    def add_picture(
        self,
        slide_index: int,
        path: str | Path,
        left: float = 1.0,
        top: float = 1.0,
        width: float | None = None,
        height: float | None = None,
    ) -> Any:
        """Insert a picture on ``slide_index`` at the given (inch) coordinates."""
        from pptx.util import Inches

        if not (0 <= slide_index < len(self.prs.slides)):
            raise IndexError(f'slide_index out of range: {slide_index}')
        slide = self.prs.slides[slide_index]
        kwargs: dict[str, Any] = {}
        if width is not None:
            kwargs['width'] = Inches(width)
        if height is not None:
            kwargs['height'] = Inches(height)
        return slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)

    # ---- Media (video / audio) ------------------------------------------- #

    def add_movie(
        self,
        slide_index: int,
        media_path: str | Path,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 6.0,
        height: float = 4.5,
        poster: str | Path | None = None,
    ) -> Any:
        """Insert a video at inch coordinates on the given slide.

        The video plays on click (autoplay timing injection is not supported
        by python-pptx and is intentionally out of scope). MP4 is the most
        interoperable container; some players reject other codecs.
        """
        from pptx.util import Inches

        if not (0 <= slide_index < len(self.prs.slides)):
            raise IndexError(f'slide_index out of range: {slide_index}')
        path = Path(media_path)
        if not path.exists():
            raise FileNotFoundError(f'media file not found: {path}')
        mime = _MEDIA_MIME.get(path.suffix.lower(), 'video/unknown')
        poster_arg = str(poster) if poster is not None else None
        slide = self.prs.slides[slide_index]
        return slide.shapes.add_movie(
            str(path),
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            poster_frame_image=poster_arg,
            mime_type=mime,
        )

    def add_audio(
        self,
        slide_index: int,
        media_path: str | Path,
        left: float = 1.0,
        top: float = 1.0,
    ) -> Any:
        """Insert an audio clip (rendered as a small speaker-shaped media shape).

        python-pptx has no dedicated audio API, so this routes through
        ``add_movie`` with an audio MIME type, which PowerPoint renders as a
        clickable audio object.
        """
        ext = Path(media_path).suffix.lower()
        if ext not in ('.mp3', '.wav', '.m4a'):
            raise ValueError(f'Unsupported audio format: {ext!r}. Use one of: .mp3, .wav, .m4a')
        return self.add_movie(slide_index, media_path, left=left, top=top, width=1.0, height=1.0)

    def add_shape(
        self,
        slide_index: int,
        shape_type: str = 'rectangle',
        left: float = 1.0,
        top: float = 1.0,
        width: float = 2.0,
        height: float = 1.0,
        fill: str | None = None,
        line: str | None = None,
    ) -> Any:
        """Add an autoshape (rectangle/oval/etc.) at the given (inch) coordinates."""
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        if not (0 <= slide_index < len(self.prs.slides)):
            raise IndexError(f'slide_index out of range: {slide_index}')
        slide = self.prs.slides[slide_index]
        shape_enum = getattr(MSO_SHAPE, shape_type.upper(), MSO_SHAPE.RECTANGLE)
        sp = slide.shapes.add_shape(
            shape_enum, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if fill:
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor.from_string(fill)
        if line:
            sp.line.color.rgb = RGBColor.from_string(line)
        return sp

    def _add_text_with_math(self, tf: Any, text: str, style: TextStyle) -> None:
        """Fill ``tf`` with ``text``, converting inline/display math to OMML."""
        import re

        from pptx.enum.text import PP_ALIGN

        has_math = '$$' in text or ('$' in text and re.search(r'\$[^$]+\$', text))
        if has_math:
            parts = re.split(r'(\$\$|\$)', text)
            p = tf.paragraphs[0]
            in_math = False
            math_display = False
            buffer = ''
            for part in parts:
                if part == '$$':
                    if in_math and math_display and buffer:
                        self._add_omath_to_paragraph(p, buffer)
                        buffer = ''
                    math_display = not math_display
                    in_math = math_display
                elif part == '$':
                    if in_math and not math_display and buffer:
                        self._add_omath_to_paragraph(p, buffer)
                        buffer = ''
                    math_display = False
                    in_math = not in_math
                elif in_math:
                    buffer += part
                else:
                    run = p.add_run()
                    run.text = part
                    self._apply_run_style(run, style)
            if buffer and in_math:
                self._add_omath_to_paragraph(p, buffer)
        else:
            tf.text = text
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    self._apply_run_style(run, style)
                if style.alignment:
                    align_map = {
                        'left': PP_ALIGN.LEFT,
                        'center': PP_ALIGN.CENTER,
                        'right': PP_ALIGN.RIGHT,
                        'justify': PP_ALIGN.JUSTIFY,
                    }
                    paragraph.alignment = align_map.get(style.alignment, PP_ALIGN.LEFT)

    def _apply_run_style(self, run: Any, style: TextStyle) -> None:
        if style.bold is not None and style.bold:
            run.font.bold = True
        if style.italic is not None and style.italic:
            run.font.italic = True
        if style.font_name is not None:
            run.font.name = style.font_name
        if style.font_size is not None:
            run.font.size = Pt(style.font_size)
        if style.color is not None:
            try:
                from pptx.dml.color import RGBColor

                run.font.color.rgb = RGBColor.from_string(style.color)
            except Exception:  # noqa: S110  # invalid color hex: leave default, continue styling
                pass

    def add_styled_content(
        self,
        tokens: list[dict[str, Any]],
    ) -> Any:
        """Render styled LaTeX-style tokens onto slides and return the last slide."""
        slide = self.add_slide()
        current_style = self._base_style

        for token in tokens:
            content_kind = token.get('type', 'text')
            content = token.get('content', '')

            if content_kind == 'text':
                if content.strip():
                    self._append_text_to_slide(slide, content, current_style)
            elif content_kind == 'command':
                cmd = token.get('command', '')
                if cmd in ('newpage',):
                    slide = self.add_slide()
                    continue
                if cmd == 'heading':
                    h_content = token.get('content', '')
                    self._append_text_to_slide(slide, h_content, current_style, heading=True)
                    continue
                if cmd == 'math':
                    latex = token.get('latex', '')
                    self._add_omath_to_slide(slide, latex, current_style)
                    continue
                token_style = TextStyle.from_latex_token(token)
                merged = current_style.merge(token_style)
                inner = token.get('content', '')
                if inner:
                    self._append_text_to_slide(slide, inner, merged)

        return slide

    def _add_omath_to_slide(self, slide: Any, latex: str, style: TextStyle) -> None:
        from tianshang_scribe.rendering.math_omml import latex_to_omml

        textbox = self._place_textbox(slide)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        omml = latex_to_omml(latex)
        if omml is not None:
            p._p.append(omml)

    def _add_omath_to_paragraph(self, paragraph: Any, latex: str) -> None:
        from tianshang_scribe.rendering.math_omml import latex_to_omml

        omml = latex_to_omml(latex)
        if omml is not None:
            paragraph._p.append(omml)

    def _place_textbox(
        self,
        slide: Any,
        left: float = 1.0,
        width: float | None = None,
        height: float = 1.0,
    ) -> Any:
        """Add a text box on ``slide``, stacking below previously placed boxes."""
        from pptx.util import Inches

        if width is None:
            slide_width_in = (self.prs.slide_width or 914400 * 10) / 914400
            width = max(1.0, slide_width_in - 2 * left)
        sid = id(slide)
        top = self._text_cursors.get(sid, 1.0)
        self._text_cursors[sid] = top + height + 0.1
        return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

    def _append_text_to_slide(
        self, slide: Any, text: str, style: TextStyle, heading: bool = False
    ) -> None:
        textbox = self._place_textbox(slide)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if heading:
            p.level = 0
        run = p.add_run()
        run.text = text
        self._apply_run_style(run, style)

    def add_latex_content(self, text: str) -> Any:
        """Parse LaTeX-style markup and render it onto slides."""
        from tianshang_scribe.rendering.latex_parser import parse_structured

        tokens = parse_structured(text)
        return self.add_styled_content(tokens)

    def replace_text(self, old: str, new: str, regex: bool = False) -> int:
        """Replace all occurrences of ``old`` across slides, preserving run styles.

        Unlike a naive per-run replace, this handles matches that span multiple
        runs (e.g. a word split across styled runs) by keeping each run's font
        and rewriting only the affected run text.
        """
        import re as _re

        count = 0
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    runs = paragraph.runs
                    if not runs:
                        current = paragraph.text
                        new_text = (
                            _re.sub(old, new, current) if regex else current.replace(old, new)
                        )
                        if new_text != current:
                            paragraph.text = new_text
                            count += 1
                        continue
                    full = ''.join(r.text for r in runs)
                    matches = (
                        list(_re.finditer(old, full))
                        if regex
                        else list(_re.finditer(_re.escape(old), full))
                    )
                    if not matches:
                        continue
                    spans = []
                    idx = 0
                    for r in runs:
                        spans.append((idx, idx + len(r.text), r))
                        idx += len(r.text)
                    for m in matches:
                        ms, me = m.start(), m.end()
                        repl = m.expand(new) if regex else new
                        covered = [s for s in spans if not (s[1] <= ms or s[0] >= me)]
                        if not covered:
                            continue
                        first = covered[0]
                        last = covered[-1]
                        prefix = full[first[0] : ms]
                        suffix = full[me : last[1]]
                        first[2].text = prefix + repl + suffix
                        for s in covered[1:]:
                            s[2].text = ''
                    count += len(matches)
        return count

    def set_style(self, style_str: str) -> None:
        """Merge a style string into the engine's base style."""
        new_style = TextStyle.from_string(style_str)
        self._base_style = self._base_style.merge(new_style)

    def apply_style_to_all(self) -> None:
        """Apply the base style to every text run in the presentation."""
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            self._apply_run_style(run, self._base_style)

    def to_pdf(self, output_path: str | Path) -> None:
        """Convert the presentation to a PDF at the given output path."""
        self.save()
        from tianshang_scribe.transform.pdf import ppt_to_pdf

        ppt_to_pdf(str(self._path), str(output_path))

    def add_slide(self, layout_index: int = 1) -> Any:
        """Add a new slide using the layout at the given index."""
        slide_layout = self.prs.slide_layouts[layout_index]
        return self.prs.slides.add_slide(slide_layout)

    def apply_layout(self, slide_index: int, layout_spec: str) -> None:
        """Apply the layout named or indexed by layout_spec to a slide."""
        try:
            layout_idx = int(layout_spec)
            slide_layout = self.prs.slide_layouts[layout_idx]
        except ValueError:
            slide_layout = None
            for lo in self.prs.slide_layouts:
                if lo.name.lower() == layout_spec.lower():
                    slide_layout = lo
                    break
            if slide_layout is None:
                raise ValueError(
                    f'Layout "{layout_spec}" not found. '
                    f'Available: {[lo.name for lo in self.prs.slide_layouts]}'
                ) from None
        if slide_index < len(self.prs.slides):
            assert slide_layout is not None  # noqa: S101  # mypy narrowing; layout validated above
            slide = self.prs.slides[slide_index]
            ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            layout_part = slide_layout.part
            r_id = slide.part.relate_to(
                layout_part,
                'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout',
            )
            csld = slide.element.find(f'{{{ns}}}cSld')
            if csld is not None:
                layout_node = csld.find(f'{{{ns}}}sldLayout')
                if layout_node is not None:
                    csld.remove(layout_node)
                layout_node = etree.SubElement(csld, f'{{{ns}}}sldLayout')
                layout_node.set(f'{{{r_ns}}}id', r_id)

    def delete_slide(self, index: int) -> None:
        """Delete the slide at the given index."""
        r_id = self.prs.slides._sldIdLst[index].attrib[
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        ]
        self.prs.part.drop_rel(r_id)
        del self.prs.slides._sldIdLst[index]

    def move_slide(self, from_index: int, to_index: int) -> None:
        """Move the slide at from_index to the given to_index."""
        slides = list(self.prs.slides._sldIdLst)
        slide_entry = slides.pop(from_index)
        slides.insert(to_index, slide_entry)
        self.prs.slides._sldIdLst[:] = slides

    def get_metadata(self) -> dict[str, str | None]:
        """Return presentation core properties as a metadata dictionary."""
        props = self.prs.core_properties
        return {
            'author': props.author,
            'title': props.title,
            'subject': props.subject,
            'category': props.category,
            'keywords': props.keywords,
            'comments': props.comments,
        }

    def set_metadata(self, **kwargs: str) -> None:
        """Set presentation core properties from keyword arguments."""
        props = self.prs.core_properties
        for key, value in kwargs.items():
            key_lower = key.lower()
            if key_lower == 'author':
                props.author = value
            elif key_lower == 'title':
                props.title = value
            elif key_lower == 'subject':
                props.subject = value
            elif key_lower == 'category':
                props.category = value
            elif key_lower == 'keywords':
                props.keywords = value
            elif key_lower == 'comments':
                props.comments = value

    def add_comment(self, text: str, slide_index: int = 0) -> None:
        """Append a comment to the notes slide of the given slide."""
        if slide_index >= len(self.prs.slides):
            raise ValueError(f'Slide index {slide_index} out of range.')
        slide = self.prs.slides[slide_index]
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        if tf.text:
            tf.text += '\n---\n'
        tf.text += text

    def add_notes(self, slide_index: int, text: str) -> None:
        """Set the speaker notes text for the given slide."""
        if slide_index >= len(self.prs.slides):
            raise ValueError(f'Slide index {slide_index} out of range.')
        notes_slide = self.prs.slides[slide_index].notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = text

    def to_images(self, output_dir: str | Path) -> list[Path]:
        """Export every slide to a PNG image.

        The deck is rendered to a PDF with LibreOffice (which captures *all*
        slides) and each PDF page is then rasterized to PNG. Rasterization uses
        PyMuPDF (``fitz``) when available, otherwise the poppler ``pdftoppm``
        binary. This replaces the old direct ``--convert-to png`` path, which
        LibreOffice limited to the first slide only.
        """
        import shutil
        import subprocess
        from pathlib import Path

        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        lo_bin = None
        lo_paths = [
            'libreoffice',
            'soffice',
            '/usr/bin/libreoffice',
            '/usr/bin/soffice',
            r'C:\Program Files\LibreOffice\program\soffice.exe',
        ]
        for p in lo_paths:
            if shutil.which(p):
                lo_bin = p
                break

        if not lo_bin:
            raise RuntimeError(
                'LibreOffice is required for slide image export. '
                'Install from https://www.libreoffice.org/download/'
            )

        self.save()
        pdf_dir = output_dir / '.pdf_tmp'
        pdf_dir.mkdir(parents=True, exist_ok=True)
        subprocess.run(  # noqa: S603  # fixed trusted binary; full path validated by exists()
            [
                lo_bin,
                '--headless',
                '--convert-to',
                'pdf',
                '--outdir',
                str(pdf_dir),
                str(self._path),
            ],
            check=True,
            capture_output=True,
        )
        pdf_path = next(pdf_dir.glob('*.pdf'), None)
        if pdf_path is None:
            shutil.rmtree(pdf_dir, ignore_errors=True)
            raise RuntimeError('LibreOffice failed to produce a PDF for image export.')
        self._pdf_to_png(pdf_path, output_dir)
        shutil.rmtree(pdf_dir, ignore_errors=True)
        return sorted(output_dir.glob('*.png'))

    def _pdf_to_png(self, pdf_path: str | Path, output_dir: str | Path) -> None:
        """Rasterize each page of ``pdf_path`` to ``slideN.png`` in ``output_dir``."""
        from pathlib import Path

        out = Path(output_dir)

        try:
            import fitz
        except ImportError:
            fitz = None

        if fitz is not None:
            doc = fitz.open(str(pdf_path))
            for i, page in enumerate(doc):
                page.get_pixmap().save(str(out / f'slide{i + 1}.png'))
            return

        import shutil
        import subprocess

        if shutil.which('pdftoppm'):
            subprocess.run(  # noqa: S603
                ['pdftoppm', '-png', '-r', '150', str(pdf_path), str(out / 'slide')],  # noqa: S607
                check=True,
                capture_output=True,
            )
            return

        raise RuntimeError(
            'PDF rasterization requires PyMuPDF or poppler. Install with: pip install pymupdf'
        )

    def set_transition(self, transition_type: str, slide_index: int | None = None) -> None:
        """Set a transition effect on one slide or the whole deck."""
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        valid_transitions = {
            'fade',
            'push',
            'wipe',
            'cover',
            'uncover',
            'dissolve',
            'random',
            'split',
            'strips',
            'blinds',
            'checker',
            'comb',
            'zoom',
            'glitter',
            'vortex',
            'ripple',
            'honeycomb',
        }
        ttype = transition_type.lower()
        if ttype not in valid_transitions:
            raise ValueError(
                f'Unsupported transition: {transition_type}. '
                f'Use one of: {", ".join(sorted(valid_transitions))}'
            )
        slides = (
            [self.prs.slides[slide_index]] if slide_index is not None else list(self.prs.slides)
        )
        for slide in slides:
            existing = slide.element.findall(f'{{{ns}}}transition')
            for el in existing:
                slide.element.remove(el)
            trans_el = etree.SubElement(slide.element, f'{{{ns}}}transition')
            etree.SubElement(trans_el, f'{{{ns}}}{ttype}')

    # ---- Master-level footer / slide number / date ------------------------ #

    @staticmethod
    def _remove_master_placeholders(sp_tree: Any, ph_type: str) -> None:
        for sp in list(sp_tree.findall(qn('p:sp'))):
            nv_pr = sp.find(qn('p:nvSpPr'))
            if nv_pr is None:
                continue
            ph = nv_pr.find(qn('p:nvPr'))
            if ph is None:
                continue
            marker = ph.find(qn('p:ph'))
            if marker is not None and marker.get('type') == ph_type:
                sp_tree.remove(sp)

    @staticmethod
    def _next_shape_id(sp_tree: Any) -> int:
        ids = [1]
        for cnv in sp_tree.iter(qn('p:cNvPr')):
            try:
                ids.append(int(cnv.get('id') or 1))
            except ValueError:
                continue
        return max(ids) + 1

    def _inject_master_placeholder(self, sp_tree: Any, ph_type: str, body_builder: Any) -> None:
        """Insert (or replace) a master-level placeholder shape on ``sp_tree``."""
        idx, name = _MASTER_PH_SPECS[ph_type]
        self._remove_master_placeholders(sp_tree, ph_type)
        sp = etree.SubElement(sp_tree, qn('p:sp'))
        nv_sp = etree.SubElement(sp, qn('p:nvSpPr'))
        cnv = etree.SubElement(nv_sp, qn('p:cNvPr'))
        cnv.set('id', str(self._next_shape_id(sp_tree)))
        cnv.set('name', name)
        cnv_sp = etree.SubElement(nv_sp, qn('p:cNvSpPr'))
        etree.SubElement(cnv_sp, qn('a:spLocks')).set('noGrp', '1')
        nv = etree.SubElement(nv_sp, qn('p:nvPr'))
        ph = etree.SubElement(nv, qn('p:ph'))
        ph.set('type', ph_type)
        ph.set('sz', 'quarter')
        ph.set('idx', idx)
        etree.SubElement(sp, qn('p:spPr'))
        tx_body = etree.SubElement(sp, qn('p:txBody'))
        etree.SubElement(tx_body, qn('a:bodyPr'))
        etree.SubElement(tx_body, qn('a:lstStyle'))
        para = etree.SubElement(tx_body, qn('a:p'))
        body_builder(para)

    @staticmethod
    def _static_text_run(para: Any, text: str) -> None:
        run = etree.SubElement(para, qn('a:r'))
        t = etree.SubElement(run, qn('a:t'))
        t.text = text

    @staticmethod
    def _auto_field(para: Any, field_type: str, fallback: str) -> None:
        import uuid

        fld = etree.SubElement(para, qn('a:fld'))
        fld.set('id', f'{{{str(uuid.uuid4()).upper()}}}')
        fld.set('type', field_type)
        t = etree.SubElement(fld, qn('a:t'))
        t.text = fallback

    def set_master_options(
        self,
        slide_number: bool = False,
        footer_text: str | None = None,
        date_visible: bool = False,
        date_text: str | None = None,
    ) -> None:
        """Configure master-level slide numbers, footers and dates deck-wide.

        Placeholder shapes are injected into every slide layout and every
        slide so PowerPoint renders them without further per-slide work.
        Calling again replaces the previous placeholders (idempotent).
        """
        if not any([slide_number, footer_text, date_visible]):
            return
        targets: list[Any] = []
        for master in self.prs.slide_masters:
            targets.extend(master.slide_layouts)
        targets.extend(self.prs.slides)
        for element in targets:
            sp_tree = element.element.find(qn('p:cSld')).find(qn('p:spTree'))

            def build_num(para: Any) -> None:
                self._auto_field(para, 'slidenum', '\u2039#\u203a')

            def build_footer(para: Any) -> None:
                self._static_text_run(para, footer_text or '')

            def build_date(para: Any) -> None:
                if date_text:
                    self._static_text_run(para, date_text)
                else:
                    self._auto_field(para, 'datetime1', '')

            if slide_number:
                self._inject_master_placeholder(sp_tree, 'sldNum', build_num)
            if footer_text is not None:
                self._inject_master_placeholder(sp_tree, 'ftr', build_footer)
            if date_visible:
                self._inject_master_placeholder(sp_tree, 'dt', build_date)

    def apply_theme(self, name: str) -> None:
        """Apply a built-in theme to every slide master's theme part.

        Only the two built-in palettes ``office`` (stock Office look) and
        ``dark`` are available; external .thmx files are out of scope. The
        rewrite touches ``theme1.xml`` only — ``clrScheme`` and
        ``fontScheme`` — so shapes with explicitly-set local fills keep their
        colors, while everything bound to theme colors/fonts re-renders.
        """
        palette = _THEME_PALETTES.get(name.strip().lower())
        if palette is None:
            valid = ', '.join(sorted(_THEME_PALETTES))
            raise ValueError(f'Unknown theme: {name!r}. Use one of: {valid}')
        ns = {'a': _DRAWINGML_NS}
        a = f'{{{_DRAWINGML_NS}}}'
        for master in self.prs.slide_masters:
            theme_part = master.part.part_related_by(
                'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'
            )
            root = etree.fromstring(theme_part.blob)

            clr_scheme = root.find('.//a:clrScheme', ns)
            if clr_scheme is not None:
                for slot in ('dk1', 'lt1', 'dk2', 'lt2', 'hlink', 'folHlink'):
                    el = clr_scheme.find(f'a:{slot}', ns)
                    if el is None:
                        continue
                    for child in list(el):
                        el.remove(child)
                    etree.SubElement(el, f'{a}srgbClr').set('val', palette[slot])
                for i in range(1, 7):
                    slot = f'accent{i}'
                    el = clr_scheme.find(f'a:{slot}', ns)
                    if el is None:
                        continue
                    for child in list(el):
                        el.remove(child)
                    etree.SubElement(el, f'{a}srgbClr').set('val', palette[slot])

            font_scheme = root.find('.//a:fontScheme', ns)
            if font_scheme is not None:
                for tag, key in (('majorFont', 'major_latin'), ('minorFont', 'minor_latin')):
                    latin = font_scheme.find(f'a:{tag}/a:latin', ns)
                    if latin is not None:
                        latin.set('typeface', palette[key])

            theme_part._blob = etree.tostring(
                root, xml_declaration=True, encoding='UTF-8', standalone=True
            )

    def set_protection(self, password: str) -> None:
        """Protect the presentation with a modify-verifier password.

        Uses the ECMA-376 agile scheme: a random 16-byte salt plus a SHA-512
        hash iterated ``spinCount`` times, stored base64-encoded (this matches
        what PowerPoint expects, unlike a plaintext password).
        """
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        pres_elem = self.prs.part._element
        existing = pres_elem.findall(f'{{{ns}}}modifyVerifier')
        for el in existing:
            pres_elem.remove(el)
        salt_b64, hash_b64 = self._modify_verifier_hash(password)
        verifier = etree.SubElement(pres_elem, f'{{{ns}}}modifyVerifier')
        verifier.set('cryptProviderType', 'rsaAES')
        verifier.set('cryptAlgorithmClass', 'hash')
        verifier.set('cryptAlgorithmType', 'typeAny')
        verifier.set('cryptAlgorithmSid', '14')
        verifier.set('cryptSpinCount', '100000')
        verifier.set('hashData', hash_b64)
        verifier.set('saltData', salt_b64)

    @staticmethod
    def _modify_verifier_hash(password: str, spin_count: int = 100000) -> tuple[str, str]:
        """Return (salt_b64, hash_b64) for a modify-verifier password."""
        salt = os.urandom(16)
        digest = hashlib.sha512(password.encode('utf-16-le') + salt).digest()
        for i in range(spin_count):
            digest = hashlib.sha512(digest + struct.pack('<I', i)).digest()
        return (
            base64.b64encode(salt).decode('ascii'),
            base64.b64encode(digest).decode('ascii'),
        )

    @staticmethod
    def verify_modify_verifier(
        password: str, salt_b64: str, hash_b64: str, spin_count: int = 100000
    ) -> bool:
        """Return True if ``password`` matches the stored verifier values."""
        import base64 as _b64

        salt = _b64.b64decode(salt_b64)
        digest = hashlib.sha512(password.encode('utf-16-le') + salt).digest()
        for i in range(spin_count):
            digest = hashlib.sha512(digest + struct.pack('<I', i)).digest()
        return _b64.b64encode(digest).decode('ascii') == hash_b64

    def unprotect(self) -> None:
        """Remove the modify-verifier protection from the presentation."""
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        pres_elem = self.prs.part._element
        existing = pres_elem.findall(f'{{{ns}}}modifyVerifier')
        for el in existing:
            pres_elem.remove(el)

    def clear_content(self) -> None:
        """Clear text from all text frames in the presentation."""
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape.text_frame.clear()

    _SHARED_REL_SUFFIXES = (
        '/slideLayout',
        '/slideMaster',
        '/theme',
        '/notesMaster',
        '/notesSlide',
    )

    def merge_workbooks(self, paths: list[str]) -> None:
        """Merge slides (with their content) from other presentations into this one.

        Each slide is deep-cloned, including embedded images/media/charts, by
        re-mapping its relationships into this presentation's package.
        """
        for p in paths:
            src = Presentation(p)
            for slide in src.slides:
                self._clone_slide(slide, src)

    def _clone_slide(self, src_slide: Any, src_prs: Any) -> Any:
        """Append a faithful copy of ``src_slide`` to this presentation."""
        # map the source slide's layout to an equivalent one inside THIS package
        src_layout = src_slide.slide_layout
        try:
            idx = list(src_prs.slide_layouts).index(src_layout)
        except ValueError:
            idx = 1
        idx = min(idx, len(self.prs.slide_layouts) - 1)
        new_slide = self.prs.slides.add_slide(self.prs.slide_layouts[idx])
        # discard layout-provided default placeholders; real shapes are copied below
        for shape in list(new_slide.shapes):
            shape._element.getparent().remove(shape._element)

        src_part = src_slide.part
        dst_part = new_slide.part
        dst_pkg = dst_part.package
        r_ns = qn('r:id').split('}', 1)[0].lstrip('{')
        part_cache: dict[int, Any] = {}

        def copy_part(target: Any, reln_type: str) -> Any:
            key = id(target)
            if key in part_cache:
                return part_cache[key]
            if reln_type.endswith('/image'):
                new_target = dst_pkg.get_or_add_image_part(io.BytesIO(target.blob))
                part_cache[key] = new_target
                return new_target
            if reln_type.endswith('/media'):
                new_target = dst_pkg.get_or_add_media_part(io.BytesIO(target.blob))
                part_cache[key] = new_target
                return new_target
            tmpl = re.sub(r'(\d+)(?=\.[^/]*$|$)', '%d', str(target.partname))
            new_target = Part(
                dst_pkg.next_partname(tmpl), target.content_type, dst_pkg, target.blob
            )
            part_cache[key] = new_target
            for sub_rel in target.rels:
                if sub_rel.is_external:
                    new_target.relate_to(sub_rel.target_ref, sub_rel.reln_type, is_external=True)
                else:
                    new_target.relate_to(
                        copy_part(sub_rel.target_part, sub_rel.reln_type), sub_rel.reln_type
                    )
            return new_target

        def remap(rel: Any) -> str:
            if rel.is_external:
                return str(dst_part.relate_to(rel.target_ref, rel.reln_type, is_external=True))
            return str(dst_part.relate_to(copy_part(rel.target_part, rel.reln_type), rel.reln_type))

        for shape in src_slide.shapes:
            new_el = copy.deepcopy(shape._element)
            drop = False
            for el in new_el.iter():
                for attr in list(el.attrib):
                    if not (attr.startswith('{') and attr.split('}', 1)[0] == r_ns):
                        continue
                    if not (
                        attr.endswith('}embed') or attr.endswith('}link') or attr == qn('r:id')
                    ):
                        continue
                    rel = src_part.rels.get(el.get(attr))
                    if rel is None or rel.reln_type.endswith(self._SHARED_REL_SUFFIXES):
                        drop = True
                        break
                    el.set(attr, remap(rel))
                if drop:
                    break
            if drop:
                continue
            new_slide.shapes._spTree.append(new_el)

        src_bg = src_slide.element.find(f'{qn("p:cSld")}/{qn("p:bg")}')
        if src_bg is not None:
            new_slide.element.find(qn('p:cSld')).append(copy.deepcopy(src_bg))
        return new_slide

    def extract_text(self) -> str:
        """Extract all slide text as plain text lines."""
        lines: list[str] = []
        for slide_idx, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = ''.join(run.text for run in paragraph.runs)
                        if text:
                            lines.append(f'[slide {slide_idx + 1}] {text}')
        return '\n'.join(lines)

    def extract_tables(self) -> list[list[list[str]]]:
        """Extract tables from all slides as lists of string tables."""
        tables: list[list[list[str]]] = []
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if getattr(shape, 'has_table', False) and shape.has_table:
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    tables.append(rows)
        return tables

    def extract_images(self, output_dir: str | Path) -> list[Path]:
        """Save picture shapes to the output dir and list the saved paths."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        out = Path(output_dir)
        out.mkdir(parents=True, exist_ok=True)
        saved: list[Path] = []
        idx = 0
        for slide_idx, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    target = out / f'slide{slide_idx + 1}_{idx}.{image.ext}'
                    target.write_bytes(image.blob)
                    saved.append(target)
                    idx += 1
        return saved

    def extract_structure(self) -> dict[str, Any]:
        """Return presentation structure summary (slide and image counts)."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        return {
            'slides': len(self.prs.slides),
            'images': sum(
                1
                for slide in self.prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ),
        }

    def compress_media(self, max_dimension: int = 1920, quality: int = 80) -> int:
        """Recompress images to reduce file size.

        Images larger than ``max_dimension`` are scaled down; JPEGs are
        re-encoded with ``quality``; PNGs are re-encoded with ``optimize``.
        Shared image parts are processed once. Returns the number of bytes
        saved (the presentation must still be saved to persist changes).
        """
        import io

        from PIL import Image as PILImage
        from pptx.oxml.ns import qn

        saved_bytes = 0
        seen: set[int] = set()
        for slide in self.prs.slides:
            for shape in slide.shapes:
                blip = shape._element.find(f'.//{qn("a:blip")}')
                if blip is None:
                    continue
                r_id = blip.get(qn('r:embed'))
                if not r_id:
                    continue
                try:
                    image_part = slide.part.related_part(r_id)
                except Exception:  # noqa: S112  # broken/missing image part: skip, do not abort export
                    continue
                part_id = id(image_part)
                if part_id in seen:
                    continue
                seen.add(part_id)

                original = image_part._blob
                try:
                    image = PILImage.open(io.BytesIO(original))
                except Exception:  # noqa: S112  # corrupt image blob: skip, do not abort compression
                    continue
                fmt = (image.format or 'JPEG').upper()
                if max(image.size) > max_dimension:
                    ratio = max_dimension / max(image.size)
                    new_size = (int(image.width * ratio), int(image.height * ratio))
                    resized = image.resize(new_size, PILImage.LANCZOS)
                else:
                    resized = image

                buf = io.BytesIO()
                if fmt == 'JPEG':
                    resized.convert('RGB').save(buf, 'JPEG', quality=quality, optimize=True)
                elif fmt == 'PNG':
                    resized.save(buf, 'PNG', optimize=True)
                else:
                    continue

                new_blob = buf.getvalue()
                if len(new_blob) < len(original):
                    image_part._blob = new_blob
                    saved_bytes += len(original) - len(new_blob)
        return saved_bytes
