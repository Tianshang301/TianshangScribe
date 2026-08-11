from __future__ import annotations

import tempfile
from pathlib import Path

import pytest

from src.core.document import (
    DocumentType,
    create_document,
    detect_document_type,
    open_document,
)
from src.core.excel_engine import ExcelEngine
from src.core.ppt_engine import PptEngine
from src.core.word_engine import WordEngine


class TestDetectDocumentType:
    def test_detect_docx(self) -> None:
        assert detect_document_type('test.docx') == DocumentType.WORD

    def test_detect_xlsx(self) -> None:
        assert detect_document_type('test.xlsx') == DocumentType.EXCEL

    def test_detect_pptx(self) -> None:
        assert detect_document_type('test.pptx') == DocumentType.PPT

    def test_detect_unknown(self) -> None:
        assert detect_document_type('test.txt') == DocumentType.UNKNOWN

    def test_detect_case_insensitive(self) -> None:
        assert detect_document_type('TEST.DOCX') == DocumentType.WORD


class TestCreateDocument:
    def test_create_word(self) -> None:
        engine = create_document(DocumentType.WORD)
        assert isinstance(engine, WordEngine)
        assert engine.doc is not None

    def test_create_excel(self) -> None:
        engine = create_document(DocumentType.EXCEL)
        assert isinstance(engine, ExcelEngine)
        assert engine.wb is not None

    def test_create_ppt(self) -> None:
        engine = create_document(DocumentType.PPT)
        assert isinstance(engine, PptEngine)
        assert engine.prs is not None

    def test_create_unknown_raises(self) -> None:
        with pytest.raises(ValueError, match='Unsupported'):
            create_document(DocumentType.UNKNOWN)


class TestWordEngine:
    @pytest.fixture
    def engine(self) -> WordEngine:
        e = WordEngine()
        e.create()
        return e

    def test_create_empty(self, engine: WordEngine) -> None:
        assert len(engine.doc.paragraphs) == 0

    def test_add_text(self, engine: WordEngine) -> None:
        engine.add_text('Hello World')
        assert engine.doc.paragraphs[0].text == 'Hello World'

    def test_add_text_with_style(self, engine: WordEngine) -> None:
        engine.add_text('Bold Text', bold=True, italic=True, font_size=14)
        paragraph = engine.doc.paragraphs[0]
        run = paragraph.runs[0]
        assert run.bold is True
        assert run.italic is True
        assert run.font.size is not None

    def test_add_text_multiple(self, engine: WordEngine) -> None:
        engine.add_text('Line 1')
        engine.add_text('Line 2')
        engine.add_text('Line 3')
        assert len(engine.doc.paragraphs) == 3

    def test_replace_text_exact(self, engine: WordEngine) -> None:
        engine.add_text('Hello World')
        count = engine.replace_text('World', 'Earth')
        assert count >= 1
        assert 'Earth' in engine.doc.paragraphs[0].text

    def test_replace_text_no_match(self, engine: WordEngine) -> None:
        engine.add_text('Hello World')
        count = engine.replace_text('Mars', 'Earth')
        assert count == 0

    def test_replace_with_regex(self, engine: WordEngine) -> None:
        engine.add_text('Hello 123 World')
        count = engine.replace_text(r'\d+', 'XXX', regex=True)
        assert count >= 1

    def test_save_and_reopen(self, engine: WordEngine) -> None:
        engine.add_text('Persist Test')
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / 'test.docx'
            engine.save(path)
            assert path.exists()

            reopened = open_document(path)
            assert isinstance(reopened, WordEngine)
            paragraphs = reopened.doc.paragraphs
            assert len(paragraphs) >= 1

    def test_open_nonexistent_raises(self) -> None:
        engine = WordEngine()
        with pytest.raises(FileNotFoundError):
            engine.open('nonexistent_file.docx')

    def test_get_metadata(self, engine: WordEngine) -> None:
        metadata = engine.get_metadata()
        assert 'author' in metadata
        assert 'title' in metadata

    def test_set_metadata(self, engine: WordEngine) -> None:
        engine.set_metadata(author='Test Author', title='Test Title')
        metadata = engine.get_metadata()
        assert metadata['author'] == 'Test Author'
        assert metadata['title'] == 'Test Title'

    def test_set_style(self, engine: WordEngine) -> None:
        engine.set_style('font=Arial,size=14')
        assert engine._base_style.font_name == 'Arial'
        assert engine._base_style.font_size == 14

    def test_get_base_style(self, engine: WordEngine) -> None:
        style = engine.get_base_style()
        assert style.font_name == 'Times New Roman'
        assert style.font_size == 12

    def test_add_text_inherits_base_style(self, engine: WordEngine) -> None:
        engine.set_style('font=Courier New,size=18,bold')
        engine.add_text('Styled text')
        paragraph = engine.doc.paragraphs[0]
        run = paragraph.runs[0]
        assert run.font.name == 'Courier New'
        assert run.bold is True

    def test_add_latex_content(self, engine: WordEngine) -> None:
        engine.add_latex_content(r'\bfseries{bold} and \itshape{italic}')
        text_content = engine.doc.paragraphs[0].text
        assert 'bold' in text_content
        assert 'italic' in text_content

    def test_add_styled_content_basic(self, engine: WordEngine) -> None:
        engine.set_style('font=Arial')
        engine.add_styled_content(
            [
                {'type': 'text', 'content': 'Hello '},
                {'type': 'command', 'command': 'bfseries', 'content': 'Bold'},
            ]
        )
        paragraph = engine.doc.paragraphs[0]
        assert 'Hello' in paragraph.text
        assert 'Bold' in paragraph.text

    def test_add_toc(self, engine: WordEngine) -> None:
        engine.add_text('Intro')
        engine.add_toc()
        full_text = '\n'.join(p.text for p in engine.doc.paragraphs)
        assert 'TOC' in full_text or len(engine.doc.paragraphs) >= 2

    def test_add_section_break(self, engine: WordEngine) -> None:
        engine.add_text('Sec1')
        engine.add_section_break()
        assert len(engine.doc.sections) >= 2

    def test_set_header(self, engine: WordEngine) -> None:
        engine.set_header('My Header')
        section = engine.doc.sections[0]
        assert section.header.paragraphs[0].text == 'My Header'

    def test_set_footer(self, engine: WordEngine) -> None:
        engine.set_footer('Page X')
        section = engine.doc.sections[0]
        assert section.footer.paragraphs[0].text == 'Page X'

    def test_add_watermark(self, engine: WordEngine) -> None:
        engine.add_watermark('DRAFT')
        section = engine.doc.sections[0]
        assert section.header.paragraphs[0].text == 'DRAFT'

    def test_small_caps_via_latex(self, engine: WordEngine) -> None:
        engine.add_latex_content(r'\scshape{Small Caps Text}')
        text = engine.doc.paragraphs[0].text
        assert 'Small Caps Text' in text


class TestExcelEngine:
    @pytest.fixture
    def engine(self) -> ExcelEngine:
        e = ExcelEngine()
        e.create()
        return e

    def test_create_empty(self, engine: ExcelEngine) -> None:
        ws = engine.wb.active
        assert ws is not None

    def test_add_text(self, engine: ExcelEngine) -> None:
        engine.add_text('Cell Content')
        ws = engine.wb.active
        assert ws.cell(row=1, column=1).value == 'Cell Content'

    def test_save_and_reopen(self, engine: ExcelEngine) -> None:
        engine.add_text('Excel Test')
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / 'test.xlsx'
            engine.save(path)
            assert path.exists()
            reopened = open_document(path)
            assert isinstance(reopened, ExcelEngine)
            ws = reopened.wb.active
            assert ws.cell(row=1, column=1).value == 'Excel Test'

    def test_add_text_with_column(self, engine: ExcelEngine) -> None:
        engine.add_text('Col2', column=2)
        ws = engine.wb.active
        assert ws.cell(row=1, column=2).value == 'Col2'

    def test_add_text_multiline(self, engine: ExcelEngine) -> None:
        engine.add_text('A\nB\nC')
        ws = engine.wb.active
        assert ws.cell(row=1, column=1).value == 'A'
        assert ws.cell(row=2, column=1).value == 'B'
        assert ws.cell(row=3, column=1).value == 'C'

    def test_replace_text_exact(self, engine: ExcelEngine) -> None:
        engine.add_text('Hello World')
        count = engine.replace_text('World', 'Earth')
        assert count >= 1
        assert engine.wb.active.cell(row=1, column=1).value == 'Hello Earth'

    def test_get_metadata(self, engine: ExcelEngine) -> None:
        metadata = engine.get_metadata()
        assert 'author' in metadata
        assert 'title' in metadata

    def test_set_metadata(self, engine: ExcelEngine) -> None:
        engine.set_metadata(author='Test Author', title='Test')
        metadata = engine.get_metadata()
        assert metadata['author'] == 'Test Author'
        assert metadata['title'] == 'Test'

    def test_set_style(self, engine: ExcelEngine) -> None:
        engine.set_style('font=Courier New,size=14')
        assert engine._base_style.font_name == 'Courier New'
        assert engine._base_style.font_size == 14

    def test_add_sheet(self, engine: ExcelEngine) -> None:
        count_before = len(engine.wb.sheetnames)
        engine.add_sheet('NewSheet')
        assert len(engine.wb.sheetnames) == count_before + 1
        assert 'NewSheet' in engine.wb.sheetnames

    def test_delete_sheet(self, engine: ExcelEngine) -> None:
        engine.add_sheet('ToDelete')
        engine.delete_sheet('ToDelete')
        assert 'ToDelete' not in engine.wb.sheetnames

    def test_rename_sheet(self, engine: ExcelEngine) -> None:
        engine.add_sheet('OldName')
        engine.rename_sheet('OldName', 'NewName')
        assert 'NewName' in engine.wb.sheetnames
        assert 'OldName' not in engine.wb.sheetnames

    def test_set_column_width(self, engine: ExcelEngine) -> None:
        engine.set_column_width(2, 25.5)
        ws = engine.wb.active
        assert ws.column_dimensions['B'].width == 25.5

    def test_set_row_height(self, engine: ExcelEngine) -> None:
        engine.set_row_height(3, 30.0)
        ws = engine.wb.active
        assert ws.row_dimensions[3].height == 30.0

    def test_set_formula(self, engine: ExcelEngine) -> None:
        engine.set_formula('A1', '=SUM(B1:B10)')
        ws = engine.wb.active
        assert ws['A1'].value == '=SUM(B1:B10)'

    def test_import_csv(self, engine: ExcelEngine) -> None:
        import csv

        with tempfile.TemporaryDirectory() as tmpdir:
            csv_path = Path(tmpdir) / 'test.csv'
            with open(csv_path, 'w', newline='') as f:
                writer = csv.writer(f)
                writer.writerow(['Name', 'Age'])
                writer.writerow(['Alice', '30'])
            engine.import_csv(str(csv_path))
            ws = engine.wb.active
            assert ws.cell(row=1, column=1).value == 'Name'
            assert ws.cell(row=2, column=2).value == '30'

    def test_export_csv(self, engine: ExcelEngine) -> None:
        engine.add_text('Name', column=1)
        engine.add_text('Alice', column=1)
        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / 'out.csv'
            engine.export_csv(str(out_path))
            assert out_path.exists()
            content = out_path.read_text()
            assert 'Name' in content
            assert 'Alice' in content

    def test_export_json(self, engine: ExcelEngine) -> None:
        engine.add_text('Name', column=1)
        engine.add_text('Alice', column=1)
        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / 'out.json'
            engine.export_json(str(out_path))
            assert out_path.exists()
            import json

            data = json.loads(out_path.read_text())
            assert isinstance(data, list)
            assert data[0]['Name'] == 'Alice'

    def test_export_html(self, engine: ExcelEngine) -> None:
        engine.add_text('Header', column=1)
        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / 'out.html'
            engine.export_html(str(out_path))
            assert out_path.exists()
            content = out_path.read_text()
            assert '<table>' in content
            assert 'Header' in content

    def test_sort_asc(self, engine: ExcelEngine) -> None:
        engine.add_text('B', column=1)
        engine.add_text('A', column=1)
        engine.add_text('C', column=1)
        engine.sort('A1:A3', 'asc')
        ws = engine.wb.active
        assert ws.cell(row=1, column=1).value == 'A'
        assert ws.cell(row=2, column=1).value == 'B'
        assert ws.cell(row=3, column=1).value == 'C'

    def test_sort_desc(self, engine: ExcelEngine) -> None:
        engine.add_text('A', column=1)
        engine.add_text('B', column=1)
        engine.add_text('C', column=1)
        engine.sort('A1:A3', 'desc')
        ws = engine.wb.active
        assert ws.cell(row=1, column=1).value == 'C'
        assert ws.cell(row=3, column=1).value == 'A'

    def test_add_chart_bar(self, engine: ExcelEngine) -> None:
        ws = engine.wb.active
        engine.add_text('Cats', column=1)
        engine.add_text('10', column=1)
        engine.add_text('Dogs', column=2)
        engine.add_text('20', column=2)
        sheet_name = ws.title
        engine.add_chart('bar', f"'{sheet_name}'!A1:B2")
        assert len(ws._charts) >= 1

    def test_add_chart_unsupported_raises(self, engine: ExcelEngine) -> None:
        with pytest.raises(ValueError, match='Unsupported chart type'):
            engine.add_chart('scatter', 'A1:B2')

    def test_merge_workbooks(self, engine: ExcelEngine) -> None:
        engine.add_sheet('Data')
        engine.add_text('Hello', column=1)
        with tempfile.TemporaryDirectory() as tmpdir:
            p1 = Path(tmpdir) / 'src.xlsx'
            engine.save(str(p1))
            e2 = ExcelEngine()
            e2.create()
            e2.merge_workbooks([str(p1)])
            assert 'Data' in e2.wb.sheetnames

    def test_split_by_sheet(self, engine: ExcelEngine) -> None:
        engine.add_sheet('Sheet1')
        engine.add_sheet('Sheet2')
        with tempfile.TemporaryDirectory() as tmpdir:
            results = engine.split_by_sheet(tmpdir)
            assert len(results) == 3
            for p in results:
                assert p.exists()

    def test_clear_content(self, engine: ExcelEngine) -> None:
        engine.add_text('Data', column=1)
        engine.add_text('More', column=1)
        engine.clear_content()
        ws = engine.wb.active
        assert ws.cell(row=1, column=1).value is None
        assert ws.cell(row=2, column=1).value is None

    def test_set_protection(self, engine: ExcelEngine) -> None:
        engine.set_protection('secret123')
        assert engine.wb.security.workbook_password is not None

    def test_unprotect(self, engine: ExcelEngine) -> None:
        engine.set_protection('secret123')
        engine.unprotect()
        assert engine.wb.security.workbook_password is not None

    def test_add_comment(self, engine: ExcelEngine) -> None:
        engine.add_text('Cell', column=1)
        engine.add_comment('A1', 'A note')
        ws = engine.wb.active
        assert ws['A1'].comment is not None
        assert ws['A1'].comment.text == 'A note'


class TestPptEngine:
    @pytest.fixture
    def engine(self) -> PptEngine:
        e = PptEngine()
        e.create()
        return e

    def test_create_has_slides(self, engine: PptEngine) -> None:
        assert len(engine.prs.slides) >= 0

    def test_add_slide(self, engine: PptEngine) -> None:
        initial = len(engine.prs.slides)
        engine.add_slide()
        assert len(engine.prs.slides) == initial + 1

    def test_save_and_reopen(self, engine: PptEngine) -> None:
        engine.add_slide()
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / 'test.pptx'
            engine.save(path)
            assert path.exists()
            reopened = open_document(path)
            assert isinstance(reopened, PptEngine)

    def test_delete_slide(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.add_slide()
        engine.add_slide()
        initial = len(engine.prs.slides)
        engine.delete_slide(0)
        assert len(engine.prs.slides) == initial - 1

    def test_move_slide(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.add_slide()
        engine.add_slide()
        engine.move_slide(0, 1)
        assert len(engine.prs.slides) == 3

    def test_add_notes(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.add_notes(0, 'Speaker note')
        notes = engine.prs.slides[0].notes_slide
        assert notes.notes_text_frame.text == 'Speaker note'

    def test_set_transition(self, engine: PptEngine) -> None:
        engine.add_slide()
        engine.set_transition('fade')
        slide = engine.prs.slides[0]
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        trans = slide.element.find(f'{{{ns}}}transition')
        assert trans is not None
        assert trans.find(f'{{{ns}}}fade') is not None

    def test_set_transition_invalid_raises(self, engine: PptEngine) -> None:
        engine.add_slide()
        with pytest.raises(ValueError, match='Unsupported transition'):
            engine.set_transition('nonexistent')

    def test_set_protection(self, engine: PptEngine) -> None:
        engine.set_protection('secret')
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        pres = engine.prs.part._element
        assert pres.find(f'{{{ns}}}modifyVerifier') is not None

    def test_unprotect(self, engine: PptEngine) -> None:
        engine.set_protection('secret')
        engine.unprotect()
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        pres = engine.prs.part._element
        assert pres.find(f'{{{ns}}}modifyVerifier') is None


class TestWordExtract:
    @pytest.fixture
    def engine(self, tmp_path: Path) -> WordEngine:
        e = WordEngine()
        e.create()
        e.add_text('Hello World')
        e.add_table_data([['Name', 'City'], ['Alice', 'NYC']])
        from PIL import Image

        img = tmp_path / 'pixel.png'
        Image.new('RGB', (32, 32), (200, 40, 40)).save(img)
        e.add_image(str(img))
        return e

    def test_extract_text_contains_paragraph(self, engine: WordEngine) -> None:
        assert 'Hello World' in engine.extract_text()

    def test_extract_text_contains_table_row(self, engine: WordEngine) -> None:
        assert 'Alice | NYC' in engine.extract_text()

    def test_extract_tables(self, engine: WordEngine) -> None:
        tables = engine.extract_tables()
        assert tables == [[['Name', 'City'], ['Alice', 'NYC']]]

    def test_extract_images(self, engine: WordEngine, tmp_path: Path) -> None:
        saved = engine.extract_images(str(tmp_path / 'imgs'))
        assert len(saved) == 1
        assert saved[0].exists()

    def test_extract_structure(self, engine: WordEngine) -> None:
        struct = engine.extract_structure()
        assert struct['tables'] == 1
        assert struct['images'] == 1
        assert struct['sections'] == 1


class TestExcelImportJson:
    def test_import_json_dicts(self, tmp_path: Path) -> None:
        data = [{'name': 'Alice', 'age': 30}, {'name': 'Bob', 'age': 25}]
        json_path = tmp_path / 'data.json'
        json_path.write_text(__import__('json').dumps(data), encoding='utf-8')
        e = ExcelEngine()
        e.create()
        e.import_json(str(json_path))
        rows = [list(r) for r in e.wb.active.iter_rows(values_only=True)]
        assert rows == [['name', 'age'], ['Alice', 30], ['Bob', 25]]

    def test_import_json_arrays(self, tmp_path: Path) -> None:
        data = [[1, 2], [3, 4]]
        json_path = tmp_path / 'data.json'
        json_path.write_text(__import__('json').dumps(data), encoding='utf-8')
        e = ExcelEngine()
        e.create()
        e.import_json(str(json_path))
        rows = [list(r) for r in e.wb.active.iter_rows(values_only=True)]
        assert rows == [[1, 2], [3, 4]]

    def test_import_json_empty_raises(self, tmp_path: Path) -> None:
        json_path = tmp_path / 'data.json'
        json_path.write_text('[]', encoding='utf-8')
        e = ExcelEngine()
        e.create()
        with pytest.raises(ValueError, match='non-empty'):
            e.import_json(str(json_path))


class TestExcelExtract:
    def test_extract_text(self) -> None:
        e = ExcelEngine()
        e.create()
        ws = e.wb.active
        ws['A1'] = 'x'
        ws['B1'] = 2
        text = e.extract_text()
        assert '[Sheet] x | 2' in text

    def test_extract_tables(self) -> None:
        e = ExcelEngine()
        e.create()
        e.wb.active['A1'] = 'a'
        tables = e.extract_tables()
        assert tables == [[['a']]]

    def test_extract_structure(self) -> None:
        e = ExcelEngine()
        e.create()
        struct = e.extract_structure()
        assert struct['sheets'] == ['Sheet']


class TestPptExtract:
    def test_extract_text(self) -> None:
        e = PptEngine()
        e.create()
        e.add_text('Slide one')
        assert '[slide 1] Slide one' in e.extract_text()

    def test_extract_structure(self) -> None:
        e = PptEngine()
        e.create()
        e.add_slide()
        struct = e.extract_structure()
        assert struct['slides'] == 1


class TestPptCompressMedia:
    def test_compress_media_reduces_size(self, tmp_path: Path) -> None:
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        img = tmp_path / 'big.jpg'
        Image.new('RGB', (4000, 2000), (128, 128, 128)).save(img, 'JPEG', quality=95)
        pptx_path = tmp_path / 'deck.pptx'
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(4))
        prs.save(str(pptx_path))
        original = pptx_path.stat().st_size

        e = PptEngine()
        e.open(str(pptx_path))
        saved = e.compress_media(max_dimension=1600, quality=60)
        assert saved > 0
        e.save(str(pptx_path))
        assert pptx_path.stat().st_size < original
