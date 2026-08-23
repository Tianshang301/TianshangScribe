"""Tests for 0.9.0 P1 master-level footer / slide number / date."""

from __future__ import annotations

from pathlib import Path

from pptx.oxml.ns import qn

from tianshang_scribe.core.ppt_engine import PptEngine


def _ph_sps(element: object, ph_type: str) -> list[object]:
    sp_tree = element.element.find(qn('p:cSld')).find(qn('p:spTree'))  # type: ignore[attr-defined]
    out: list[object] = []
    for sp in sp_tree.findall(qn('p:sp')):
        nv_pr = sp.find(qn('p:nvSpPr'))
        if nv_pr is None:
            continue
        ph = nv_pr.find(qn('p:nvPr'))
        if ph is None:
            continue
        marker = ph.find(qn('p:ph'))
        if marker is not None and marker.get('type') == ph_type:
            out.append(sp)
    return out


def _sp_text(sp: object) -> str:
    return ''.join(t.text or '' for t in sp.iter(qn('a:t')))


def _deck() -> PptEngine:
    e = PptEngine()
    e.create()
    e.add_slide()
    e.add_slide()
    return e


class TestSetMasterOptions:
    def test_full_configuration_persists(self, tmp_path: Path) -> None:
        e = _deck()
        e.set_master_options(
            slide_number=True,
            footer_text='Confidential',
            date_visible=True,
            date_text='2026-08-23',
        )
        out = tmp_path / 'master.pptx'
        e.save(out)
        reopened = PptEngine()
        reopened.open(out)
        for slide in reopened.prs.slides:
            footers = _ph_sps(slide, 'ftr')
            assert len(footers) == 1
            assert _sp_text(footers[0]) == 'Confidential'
            nums = _ph_sps(slide, 'sldNum')
            assert len(nums) == 1
            assert _sp_text(nums[0]) == '\u2039#\u203a'
            dates = _ph_sps(slide, 'dt')
            assert len(dates) == 1
            assert _sp_text(dates[0]) == '2026-08-23'
        # layouts carry the placeholders too
        for master in reopened.prs.slide_masters:
            for layout in master.slide_layouts:
                assert len(_ph_sps(layout, 'ftr')) == 1

    def test_repeated_call_is_idempotent(self, tmp_path: Path) -> None:
        e = _deck()
        e.set_master_options(footer_text='v1')
        e.set_master_options(footer_text='v2')
        out = tmp_path / 'idem.pptx'
        e.save(out)
        reopened = PptEngine()
        reopened.open(out)
        for slide in reopened.prs.slides:
            footers = _ph_sps(slide, 'ftr')
            assert len(footers) == 1
            assert _sp_text(footers[0]) == 'v2'

    def test_noop_configuration_injects_nothing(self, tmp_path: Path) -> None:
        e = _deck()
        e.set_master_options()
        out = tmp_path / 'noop.pptx'
        e.save(out)
        reopened = PptEngine()
        reopened.open(out)
        for slide in reopened.prs.slides:
            assert _ph_sps(slide, 'ftr') == []
            assert _ph_sps(slide, 'sldNum') == []
            assert _ph_sps(slide, 'dt') == []
